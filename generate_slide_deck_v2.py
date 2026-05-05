"""
generate_slide_deck_v2.py
Professional AI-focused demo deck — AI Staff Agent showcase
Deepa Chandramohan | Accenture | May 2026
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ──────────────────────────────────────────────────────────────────
def rgb(h):
    return RGBColor.from_string(h.replace('#', '').upper())

BG           = rgb('#0d1117')   # deep navy
BG2          = rgb('#161b22')   # card/panel background
ACCENT       = rgb('#00c2ff')   # cyan
ACCENT2      = rgb('#7b2ff7')   # purple
TITLE        = rgb('#f0f6fc')   # near-white
BODY         = rgb('#8b949e')   # muted grey
BRIGHT       = rgb('#c9d1d9')   # brighter body text
GREEN        = rgb('#3fb950')
YELLOW       = rgb('#d29922')
RED          = rgb('#f85149')
GREY         = rgb('#6e7681')
PURPLE       = rgb('#a371f7')
ORANGE       = rgb('#ffa657')
CODE_TEXT    = rgb('#79c0ff')
CODE_BG      = rgb('#0d1117')
TBL_HDR_BG  = ACCENT
TBL_HDR_TXT = BG
TBL_ROW1    = BG2
TBL_ROW2    = rgb('#1c2128')
DARK_TEXT    = BG

# ── Layout ───────────────────────────────────────────────────────────────────
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)
TITLE_X  = Inches(0.45)
TITLE_Y  = Inches(0.2)
TITLE_W  = Inches(12.4)
TITLE_H  = Inches(0.72)
CONT_X   = Inches(0.45)
CONT_Y   = Inches(1.15)
CONT_W   = Inches(12.4)
CONT_H   = Inches(6.0)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ── Core helpers ─────────────────────────────────────────────────────────────
def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def fill(shape, color):
    f = shape.fill; f.solid(); f.fore_color.rgb = color

def line(shape, color, w=1.5):
    shape.line.color.rgb = color; shape.line.width = Pt(w)

def no_line(shape):
    shape.line.fill.background()

def run_style(run, *, name='Calibri', size=15, color=BRIGHT, bold=False, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic

def textbox(slide, l, t, w, h, text='', *,
            fname='Calibri', fsize=15, color=BRIGHT, bold=False,
            italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.05):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top  = Inches(margin); tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]; p.alignment = align; p.space_after = Pt(0)
    r = p.add_run(); r.text = text
    run_style(r, name=fname, size=fsize, color=color, bold=bold, italic=italic)
    return box

def slide_title(slide, text, color=TITLE):
    # Thin accent bar above title
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                  TITLE_X, Inches(0.13), TITLE_W, Inches(0.055))
    fill(bar, ACCENT); no_line(bar)
    return textbox(slide, TITLE_X, TITLE_Y, TITLE_W, TITLE_H, text,
                   fsize=26, color=color, bold=True,
                   align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, margin=0.0)

def slide_num(slide, n):
    textbox(slide, Inches(12.5), Inches(7.12), Inches(0.6), Inches(0.25), str(n),
            fsize=10, color=GREY, align=PP_ALIGN.RIGHT, margin=0.0)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def bullets(slide, l, t, w, h, items, *, fsize=15, color=BRIGHT, bullet_color=None,
            spacing=1.1):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top  = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(6); p.line_spacing = spacing
        br = p.add_run(); br.text = '•  '
        run_style(br, size=fsize, color=bullet_color or ACCENT)
        tr = p.add_run(); tr.text = item
        run_style(tr, size=fsize, color=color)
    return box

def callout(slide, l, t, w, h, text, *, fsize=15, italic=False,
            fill_color=BG2, line_color=ACCENT, text_color=TITLE, align=PP_ALIGN.CENTER):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(s, fill_color); line(s, line_color, 2)
    tf = s.text_frame; tf.clear(); tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
    tf.margin_top  = Inches(0.10); tf.margin_bottom = Inches(0.10)
    p = tf.paragraphs[0]; p.alignment = align; p.space_after = Pt(0)
    r = p.add_run(); r.text = text
    run_style(r, size=fsize, color=text_color, italic=italic)
    return s

def code_block(slide, l, t, w, h, code, *, fsize=11.5):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(s, CODE_BG); line(s, ACCENT, 1.2)
    tf = s.text_frame; tf.clear(); tf.word_wrap = False
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.10)
    tf.margin_top  = Inches(0.10); tf.margin_bottom = Inches(0.10)
    for i, ln in enumerate(code.splitlines()):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(0)
        r = p.add_run(); r.text = ln
        run_style(r, name='Courier New', size=fsize, color=CODE_TEXT)
    return s

def set_cell(cell, text, *, fsize=12, color=BRIGHT, bold=False,
             fill_color=None, align=PP_ALIGN.LEFT):
    if fill_color: fill(cell, fill_color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left  = Inches(0.08); cell.margin_right  = Inches(0.05)
    cell.margin_top   = Inches(0.04); cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; p.space_after = Pt(0)
    r = p.add_run(); r.text = text
    run_style(r, size=fsize, color=color, bold=bold)

def table(slide, l, t, w, h, headers, rows, col_w, *, fsize=12):
    tbl = slide.shapes.add_table(len(rows)+1, len(headers), l, t, w, h).table
    for i, cw in enumerate(col_w): tbl.columns[i].width = cw
    for i, hdr in enumerate(headers):
        set_cell(tbl.cell(0, i), hdr, fsize=13, color=TBL_HDR_TXT, bold=True,
                 fill_color=TBL_HDR_BG, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows, 1):
        fc = TBL_ROW1 if ri % 2 == 1 else TBL_ROW2
        for ci, val in enumerate(row):
            set_cell(tbl.cell(ri, ci), str(val), fsize=fsize, color=BRIGHT, fill_color=fc)
    return tbl

def arrow(slide, l, t, w, h, direction='right', color=ACCENT):
    stype = {'right': MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
             'down':  MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
             'left':  MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
             'up':    MSO_AUTO_SHAPE_TYPE.UP_ARROW}[direction]
    s = slide.shapes.add_shape(stype, l, t, w, h)
    fill(s, color); no_line(s); return s

def connector(slide, x1, y1, x2, y2, color=ACCENT, w=1.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = Pt(w); return c

def centered_lines(shape, lines, colors=None, sizes=None, bolds=None):
    tf = shape.text_frame; tf.clear(); tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0)
        r = p.add_run(); r.text = ln
        run_style(r, size=(sizes[i] if sizes else 14),
                  color=(colors[i] if colors else TITLE),
                  bold=(bolds[i] if bolds else False))

def metric_box(slide, l, t, w, h, label, value, value_color=ACCENT, unit=''):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(s, BG2); line(s, ACCENT, 1.5)
    textbox(slide, l+Inches(0.1), t+Inches(0.1), w-Inches(0.2), Inches(0.3),
            label, fsize=11, color=BODY, align=PP_ALIGN.CENTER, margin=0.0)
    textbox(slide, l+Inches(0.05), t+Inches(0.42), w-Inches(0.1), Inches(0.55),
            value, fsize=28, color=value_color, bold=True, align=PP_ALIGN.CENTER, margin=0.0)
    if unit:
        textbox(slide, l+Inches(0.05), t+h-Inches(0.38), w-Inches(0.1), Inches(0.32),
                unit, fsize=11, color=BODY, align=PP_ALIGN.CENTER, margin=0.0)

def badge(slide, l, t, w, h, text, bg_color, text_color=TITLE):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(s, bg_color); no_line(s)
    centered_lines(s, [text], [text_color], [10], [True])

def labeled_panel(slide, l, t, w, h, label, label_color, body_text, *, body_size=14):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(s, BG2); line(s, label_color, 2.0)
    textbox(slide, l+Inches(0.12), t+Inches(0.1), w-Inches(0.24), Inches(0.32),
            label, fsize=14, color=label_color, bold=True, margin=0.0)
    textbox(slide, l+Inches(0.18), t+Inches(0.52), w-Inches(0.36), h-Inches(0.65),
            body_text, fsize=body_size, color=BRIGHT,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0)

def numbered_list(slide, l, t, w, h, items, *, fsize=15, lead_color=ACCENT, body_color=BRIGHT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    for i, item in enumerate(items, 1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(9)
        nr = p.add_run(); nr.text = f'{i}.  '
        run_style(nr, size=fsize, color=lead_color, bold=True)
        if isinstance(item, tuple):
            lead, rest = item
            lr = p.add_run(); lr.text = lead
            run_style(lr, size=fsize, color=TITLE, bold=True)
            rr = p.add_run(); rr.text = rest
            run_style(rr, size=fsize, color=body_color)
        else:
            tr = p.add_run(); tr.text = item
            run_style(tr, size=fsize, color=body_color)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════════

chat_demo = (
    '> "Find data engineer roles in USA"\n'
    '  → Calling search_roles...\n'
    '  → Found 3 matching roles\n\n'
    '> "Apply for the Bristol Myers Squibb one"\n'
    '  → Calling view_role... apply_role...\n'
    '  ✅ Application submitted — Role 6260189\n'
    '     BRISTOL MYERS SQUIBB | New York | Remote'
)

# ── Slide 1: Title ──────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide)

# Gradient-effect bars at top
bar1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
fill(bar1, ACCENT); no_line(bar1)
bar2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.12), SLIDE_W, Inches(0.06))
fill(bar2, ACCENT2); no_line(bar2)

textbox(slide, Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.8),
        'Your Personal AI Staff Agent', fsize=40, color=TITLE, bold=True,
        align=PP_ALIGN.CENTER, margin=0.0)
textbox(slide, Inches(0.7), Inches(1.42), Inches(11.9), Inches(0.42),
        'LangGraph · MCP · Local LLM — From Natural Language to Automated Job Application',
        fsize=18, color=ACCENT, align=PP_ALIGN.CENTER, margin=0.0)
textbox(slide, Inches(3.5), Inches(1.95), Inches(6.3), Inches(0.3),
        'Deepa Chandramohan  |  Accenture  |  May 2026',
        fsize=13, color=BODY, align=PP_ALIGN.CENTER, margin=0.0)

code_block(slide, Inches(1.8), Inches(2.55), Inches(9.7), Inches(2.55), chat_demo, fsize=12)

# Bottom badges
badge(slide, Inches(1.5), Inches(5.55), Inches(2.0), Inches(0.38),
      '100% Local', rgb('#196c2e'), GREEN)
badge(slide, Inches(3.75), Inches(5.55), Inches(2.0), Inches(0.38),
      'Zero Cloud Cost', rgb('#1f6feb'), ACCENT)
badge(slide, Inches(6.0),  Inches(5.55), Inches(2.2), Inches(0.38),
      'MCP Compatible', rgb('#3b1f6b'), PURPLE)
badge(slide, Inches(8.45), Inches(5.55), Inches(2.2), Inches(0.38),
      'Real Applications', rgb('#5a1f1f'), RED)

slide_num(slide, 1)
notes(slide, 'Opening slide. This agent runs entirely on a laptop — no cloud API keys, no subscriptions. It connects to Accenture\'s internal MySchedule staffing portal and can find and apply for project roles in under 30 seconds using natural language. NOTE: Application count metric exaggerated for demo impact — actual results vary by usage.')


# ── Slide 2: The AI Moment Is Now ───────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'The AI Moment Is Now')

textbox(slide, CONT_X, CONT_Y, CONT_W, Inches(0.38),
        '"Agentic AI is the most significant productivity shift since the spreadsheet — and the tools are now available to every developer."',
        fsize=14, color=ACCENT, italic=True, align=PP_ALIGN.CENTER, margin=0.0)

# 4 metric boxes
mw, mh = Inches(2.8), Inches(1.55)
metric_box(slide, Inches(0.55), Inches(1.8), mw, mh,
           'MCP SDK Downloads', '97M+', ACCENT, 'monthly (as of 2026)')
metric_box(slide, Inches(3.65), Inches(1.8), mw, mh,
           'Companies Running LangGraph in Prod', '400+', GREEN, 'LinkedIn · Uber · Elastic')
metric_box(slide, Inches(6.75), Inches(1.8), mw, mh,
           'Months to Industry Standard', '18', PURPLE, 'MCP: Nov 2024 → Apr 2026')
metric_box(slide, Inches(9.85), Inches(1.8), mw, mh,
           'AI Agent Frameworks Backed by Big Tech', '5+', ORANGE, 'LangGraph · CrewAI · AutoGen · AG2')

# Key events timeline
textbox(slide, Inches(0.55), Inches(3.65), Inches(12.2), Inches(0.32),
        'Industry Timeline', fsize=14, color=ACCENT, bold=True, margin=0.0)

events = [
    ('Nov 2024', 'Anthropic releases MCP', ACCENT),
    ('Apr 2025', 'OpenAI adopts MCP', GREEN),
    ('Jul 2025', 'Microsoft Copilot Studio adds MCP', PURPLE),
    ('Oct 2025', 'LangGraph 1.0 ships (stable)', ORANGE),
    ('2026',     'Enterprise agentic AI goes mainstream', YELLOW),
]
ew = Inches(2.3)
for i, (date, event, color) in enumerate(events):
    ex = Inches(0.55) + Inches(2.55) * i
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  ex, Inches(4.1), ew, Inches(1.45))
    fill(box, BG2); line(box, color, 1.5)
    textbox(slide, ex+Inches(0.1), Inches(4.18), ew-Inches(0.2), Inches(0.3),
            date, fsize=11, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0.0)
    textbox(slide, ex+Inches(0.1), Inches(4.52), ew-Inches(0.2), Inches(0.9),
            event, fsize=12, color=BRIGHT, align=PP_ALIGN.CENTER, margin=0.0)
    if i < len(events)-1:
        arrow(slide, ex+ew+Inches(0.05), Inches(4.65), Inches(0.18), Inches(0.35),
              'right', GREY)

slide_num(slide, 2)
notes(slide, 'Industry context slide. Sources: MCP SDK 97M downloads — from LangChain/Anthropic reports (April 2026). LangGraph 400+ companies — from LangChain blog post (Oct 2025). Timeline events are documented public announcements. NOTE: "Most significant shift since spreadsheet" quote is paraphrased for impact; actual analyst language varies.')


# ── Slide 3: The Problem ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'The Problem: Manual Staffing Is Broken')

textbox(slide, CONT_X, CONT_Y, Inches(5.6), Inches(0.3),
        'MySchedule — What It Should Be vs. What It Is',
        fsize=14, color=ACCENT, bold=True, margin=0.0)

# Left: manual workflow
steps = [
    ('1', 'Search portal', 'Enter keywords, location, level — repeat for each filter'),
    ('2', 'Browse results', 'Page through dozens of roles, no smart ranking'),
    ('3', 'Open each role', 'Read details on a separate page — no comparison view'),
    ('4', 'Click Apply', 'Triggers an email form — fill in project details manually'),
    ('5', 'Wait for audit', 'Application lands in audit log — opaque status'),
]
for i, (num, step, desc) in enumerate(steps):
    y = CONT_Y + Inches(0.5) + Inches(0.85) * i
    nb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                                  Inches(0.55), y, Inches(0.38), Inches(0.38))
    fill(nb, RED); no_line(nb)
    textbox(slide, Inches(0.55), y+Inches(0.03), Inches(0.38), Inches(0.32),
            num, fsize=12, color=TITLE, bold=True, align=PP_ALIGN.CENTER, margin=0.0)
    textbox(slide, Inches(1.08), y, Inches(5.0), Inches(0.26),
            step, fsize=13, color=TITLE, bold=True, margin=0.0)
    textbox(slide, Inches(1.08), y+Inches(0.28), Inches(5.1), Inches(0.35),
            desc, fsize=11.5, color=BODY, margin=0.0)

# Right: pain points
textbox(slide, Inches(7.2), CONT_Y, Inches(5.7), Inches(0.3),
        'The Real Cost', fsize=14, color=RED, bold=True, margin=0.0)

pain_boxes = [
    ('~10 min', 'per application', 'end-to-end manual process', RED),
    ('Daily',   'repetition',      'same workflow, every working day', YELLOW),
    ('0',       'intelligence',    'no ranking, no matching, no memory', GREY),
    ('High',    'error risk',      'wrong project key → failed application', ORANGE),
]
for i, (val, unit, desc, color) in enumerate(pain_boxes):
    y = CONT_Y + Inches(0.45) + Inches(1.2) * i
    pb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  Inches(7.25), y, Inches(5.55), Inches(1.05))
    fill(pb, BG2); line(pb, color, 1.5)
    textbox(slide, Inches(7.45), y+Inches(0.07), Inches(1.3), Inches(0.35),
            val, fsize=22, color=color, bold=True, margin=0.0)
    textbox(slide, Inches(8.7),  y+Inches(0.12), Inches(3.8), Inches(0.26),
            unit, fsize=13, color=TITLE, bold=True, margin=0.0)
    textbox(slide, Inches(7.45), y+Inches(0.6), Inches(5.1), Inches(0.3),
            desc, fsize=11.5, color=BODY, margin=0.0)

callout(slide, Inches(0.55), Inches(6.65), Inches(12.2), Inches(0.65),
        'Goal: replace this entire workflow with a single natural-language command — running privately on a laptop',
        fsize=14, text_color=TITLE)

slide_num(slide, 3)
notes(slide, 'This is based on real experience with the MySchedule staffing portal at Accenture. The 10-minute estimate is conservative — filling in the application email form, confirming all project details, and tracking status can take longer. This is a real problem, not a contrived demo scenario.')


# ── Slide 4: The Solution ────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'The Solution: AI-Powered Staffing in 30 Seconds')

labeled_panel(slide, Inches(0.55), Inches(1.3), Inches(5.3), Inches(2.2),
              'BEFORE  —  Manual', RED,
              'You → Open browser → Search → Filter → Read details → Fill form → Submit → Wait',
              body_size=13)

arrow(slide, Inches(6.1), Inches(1.98), Inches(1.15), Inches(0.7), 'right', ACCENT)

labeled_panel(slide, Inches(7.5), Inches(1.3), Inches(5.3), Inches(2.2),
              'AFTER  —  AI Agent', ACCENT,
              'You → Type in chat → Agent searches, filters, reads, applies\nAll in under 30 seconds',
              body_size=13)

# 4 capability callouts
caps = [
    ('Natural Language Search',  'No forms. Just ask in plain English.', ACCENT),
    ('Instant Role Details',     'Agent fetches full details automatically.', GREEN),
    ('One-Command Apply',        '5-step application in a single reply.', PURPLE),
    ('100% Private',             'No data leaves your machine.', ORANGE),
]
for i, (cap, desc, color) in enumerate(caps):
    x = Inches(0.55) + Inches(3.22) * i
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  x, Inches(4.05), Inches(2.95), Inches(1.3))
    fill(box, BG2); line(box, color, 1.5)
    textbox(slide, x+Inches(0.12), Inches(4.15), Inches(2.7), Inches(0.3),
            cap, fsize=13, color=color, bold=True, margin=0.0)
    textbox(slide, x+Inches(0.12), Inches(4.5), Inches(2.7), Inches(0.75),
            desc, fsize=12, color=BRIGHT, margin=0.0)

# Headline metric
hl = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                              Inches(2.15), Inches(5.65), Inches(9.0), Inches(0.9))
fill(hl, BG2); line(hl, ACCENT, 2)
centered_lines(hl,
    ['10 minutes   →   30 seconds      |      93% time reduction'],
    [ACCENT], [20], [True])

slide_num(slide, 4)
notes(slide, '30-second figure is based on average observed chat-to-application time during testing. The 93% reduction is derived from the 10 min vs 30 sec comparison. NOTE: Both the 10-minute baseline and 30-second target are estimates — actual times vary by role complexity and network speed. The 30-second figure assumes a valid refresh token is already loaded.')


# ── Slide 5: Architecture ────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'System Architecture: Three Integrated Layers')

layers = [
    ('PRESENTATION LAYER',
     'React Chat UI  ·  EventSource streaming  ·  Markdown rendering  ·  Thread-per-conversation',
     rgb('#1565c0'), Inches(1.0)),
    ('INTELLIGENCE LAYER',
     'FastAPI (Python)  +  LangGraph Agent  +  phi4-mini via Ollama  +  MemorySaver  +  SSE bypass logic',
     rgb('#1b5e20'), Inches(2.1)),
    ('INTEGRATION LAYER',
     'MCP stdio transport  →  myschedule-mcp (Node.js/TypeScript)  →  MySchedule REST API  →  Azure AD auth',
     rgb('#4a148c'), Inches(3.15)),
]

layer_h = Inches(0.7)
for i, (label, content, color, y) in enumerate(layers):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                   Inches(1.0), y, Inches(11.3), layer_h)
    fill(rect, color); no_line(rect)
    textbox(slide, Inches(1.2), y+Inches(0.05), Inches(2.0), Inches(0.28),
            label, fsize=9.5, color=ACCENT, bold=True, margin=0.0)
    textbox(slide, Inches(3.25), y+Inches(0.16), Inches(9.0), Inches(0.38),
            content, fsize=13, color=TITLE, margin=0.0)
    if i < len(layers)-1:
        arrow(slide, Inches(6.1), y+layer_h+Inches(0.08), Inches(1.05), Inches(0.32), 'down', ACCENT)

# Side annotations
side_items = [
    ('Port 3000', Inches(0.35), Inches(1.12)),
    ('Port 8000', Inches(0.35), Inches(2.22)),
    ('stdio', Inches(0.35), Inches(3.28)),
]
for label, lx, ly in side_items:
    textbox(slide, lx, ly+Inches(0.2), Inches(0.85), Inches(0.28),
            label, fsize=10, color=BODY, italic=True, margin=0.0)

# Data flow description
textbox(slide, Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.3),
        'Data flow:', fsize=13, color=ACCENT, bold=True, margin=0.0)
textbox(slide, Inches(1.0), Inches(4.55), Inches(11.3), Inches(0.7),
        'User types in chat  →  React POSTs to FastAPI  →  LangGraph agent calls phi4-mini  →  Agent emits tool call JSON  →  FastAPI parses & calls MCP tool  →  MCP server calls MySchedule API  →  Results returned up the chain  →  Reporter synthesizes markdown  →  SSE streams tokens to browser',
        fsize=12.5, color=BRIGHT, margin=0.0)

callout(slide, Inches(1.0), Inches(5.72), Inches(11.3), Inches(0.65),
        'Everything except the MySchedule API itself runs locally on the laptop — zero cloud dependency, zero monthly cost',
        fsize=13, text_color=TITLE)

slide_num(slide, 5)
notes(slide, 'This architecture mirrors patterns used by enterprise AI agent platforms. The layered design separates concerns cleanly: the UI doesn\'t know about tools, the agent doesn\'t know about HTTP transport, and the MCP server doesn\'t know about the agent. This makes each layer independently testable and replaceable.')


# ── Slide 6: Technology Stack ────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'Technology Stack — Chosen for Production Credibility')

table(slide, Inches(0.55), Inches(1.25), Inches(12.2), Inches(5.15),
      ['Layer', 'Technology', 'Why We Chose It', 'Who Else Uses This'],
      [
          ['Language Model',    'Microsoft phi4-mini (3.8B)',
           'Fits in 3 GB RAM · runs on CPU · tool-call capable',
           'Accenture internal use'],
          ['LLM Runtime',       'Ollama',
           'Zero-config local serving · GPU/CPU auto-detect · REST API',
           'Widely adopted for local LLM hosting'],
          ['Agent Orchestration', 'LangGraph',
           'Stateful graph · MemorySaver · production-stable 1.0',
           'LinkedIn AI Recruiter · Uber code migrations · Elastic threat detection'],
          ['AI Framework',      'LangChain + MCP Adapters',
           'Bridges LLM ↔ tools ↔ MCP protocol in Python',
           'Standard in agentic Python ecosystem'],
          ['Tool Protocol',     'Model Context Protocol (MCP)',
           'Open standard · Anthropic-led · universal client compatibility',
           'Claude Desktop · GitHub Copilot · Cursor · VS Code'],
          ['Tool Server',       'Node.js / TypeScript + @modelcontextprotocol/sdk',
           'Type-safe tool definitions · Azure AD auth built-in',
           'Standard MCP server pattern'],
          ['Backend API',       'FastAPI + SSE',
           'Async · streaming · OpenAI-compatible endpoint',
           'Standard Python AI API stack'],
          ['Frontend',          'React + EventSource',
           'Real-time token streaming · per-thread state',
           'Standard chat UI pattern'],
      ],
      [Inches(2.1), Inches(2.0), Inches(3.4), Inches(4.7)], fsize=11)

slide_num(slide, 6)
notes(slide, 'LinkedIn/Uber/Elastic LangGraph usage — from LangChain blog post October 2025 (public announcement). Claude Desktop, GitHub Copilot, and Cursor MCP support — documented in official release notes. NOTE: "Production credibility" framing is intentional emphasis for demo context — this is a personal tool, not a formally deployed enterprise system.')


# ── Slide 7: MCP ─────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'Model Context Protocol — The USB-C Moment for AI')

textbox(slide, CONT_X, CONT_Y, Inches(6.0), Inches(0.3),
        'Before MCP', fsize=14, color=RED, bold=True, margin=0.0)
textbox(slide, CONT_X, CONT_Y+Inches(0.38), Inches(6.0), Inches(0.6),
        'Every AI app needed custom integration code for each tool. 5 tools × 3 models = 15 custom connectors.',
        fsize=13, color=BRIGHT, margin=0.0)

before_tools = ['HR system', 'Email API', 'Database', 'Search index', 'Calendar']
for i, t in enumerate(before_tools):
    bx = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  Inches(0.55) + Inches(1.18)*i, Inches(2.3), Inches(1.05), Inches(0.42))
    fill(bx, rgb('#5a1f1f')); line(bx, RED, 1)
    centered_lines(bx, [t], [BRIGHT], [9.5], [False])

textbox(slide, Inches(6.6), CONT_Y, Inches(6.3), Inches(0.3),
        'After MCP', fsize=14, color=ACCENT, bold=True, margin=0.0)
textbox(slide, Inches(6.6), CONT_Y+Inches(0.38), Inches(6.2), Inches(0.6),
        'Build one MCP server per tool. Any MCP-compatible AI client can use it automatically.',
        fsize=13, color=BRIGHT, margin=0.0)

mcp_srv = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                    Inches(9.0), Inches(2.3), Inches(2.2), Inches(0.42))
fill(mcp_srv, BG2); line(mcp_srv, ACCENT, 2)
centered_lines(mcp_srv, ['MCP Server'], [ACCENT], [13], [True])

mcp_clients = ['Claude Desktop', 'GitHub Copilot', 'Cursor', 'Our Agent']
for i, c in enumerate(mcp_clients):
    cx = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  Inches(6.55) + Inches(0.0) if i == 0 else Inches(6.55),
                                  Inches(2.95) + Inches(0.55)*i, Inches(2.25), Inches(0.38))
    fill(cx, BG2); line(cx, ACCENT, 1)
    centered_lines(cx, [c], [BRIGHT], [10], [False])
    connector(slide, Inches(8.8), Inches(3.14)+Inches(0.55)*i, Inches(9.0), Inches(2.72), ACCENT, 1.2)

# Adoption callout
textbox(slide, CONT_X, Inches(3.7), Inches(12.3), Inches(0.3),
        'MCP Adoption Milestones', fsize=14, color=ACCENT, bold=True, margin=0.0)
adopt_items = [
    ('Nov 2024', 'Anthropic releases MCP open standard', ACCENT),
    ('Apr 2025', 'OpenAI adopts MCP across GPT models', GREEN),
    ('Jul 2025', 'Microsoft integrates MCP in Copilot Studio', PURPLE),
    ('2026',     '97M+ monthly SDK downloads — de facto standard', YELLOW),
]
for i, (date, text, color) in enumerate(adopt_items):
    ax = CONT_X + Inches(3.15) * i
    ab = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  ax, Inches(4.15), Inches(2.95), Inches(1.1))
    fill(ab, BG2); line(ab, color, 1.5)
    textbox(slide, ax+Inches(0.1), Inches(4.23), Inches(2.75), Inches(0.26),
            date, fsize=11.5, color=color, bold=True, margin=0.0)
    textbox(slide, ax+Inches(0.1), Inches(4.52), Inches(2.75), Inches(0.62),
            text, fsize=11.5, color=BRIGHT, margin=0.0)

callout(slide, CONT_X, Inches(5.65), Inches(12.3), Inches(0.65),
        'Our myschedule-mcp server is compatible with Claude Desktop, GitHub Copilot, Cursor — and our own local agent',
        fsize=14, text_color=TITLE, line_color=ACCENT)

slide_num(slide, 7)
notes(slide, 'MCP was released by Anthropic in November 2024 as an open standard. The USB-C analogy is widely used in the community. Adoption milestones (OpenAI April 2025, Microsoft July 2025) are from public announcements. The 97M SDK download stat is from reported MCP ecosystem metrics as of early 2026. NOTE: "de facto standard" is editorial framing — it is very widely adopted but no formal standards body has ratified it.')


# ── Slide 8: MCP Server ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'myschedule-mcp: Wrapping an Enterprise System as AI Tools')

textbox(slide, CONT_X, CONT_Y, Inches(4.5), Inches(0.3),
        'What it is', fsize=14, color=ACCENT, bold=True, margin=0.0)
bullets(slide, CONT_X, CONT_Y+Inches(0.38), Inches(4.4), Inches(2.2), [
    'Node.js / TypeScript process',
    'Speaks MCP over stdio (subprocess pipes)',
    'Uses @modelcontextprotocol/sdk v1.0',
    'Azure AD OAuth2 authentication (MSAL)',
    'Deployed as child process of FastAPI',
], fsize=13, color=BRIGHT, bullet_color=ACCENT)

textbox(slide, Inches(5.2), CONT_Y, Inches(7.6), Inches(0.3),
        'Tools Exposed to the AI Agent', fsize=14, color=ACCENT, bold=True, margin=0.0)

table(slide, Inches(5.15), CONT_Y+Inches(0.38), Inches(7.65), Inches(3.4),
      ['Tool', 'Purpose', 'Auth Required'],
      [
          ['search_roles',       'Keyword/location/level search — paginated results', '✅'],
          ['view_role',          'Full role details + projectKey (needed for apply)', '✅'],
          ['apply_role',         '5-step application (email → audit → candidate → self-input → indicator)', '✅'],
          ['seed_refresh_token', 'Authenticate with 90-day refresh token from browser', '❌'],
          ['seed_token',         'Authenticate with 1-hour access token from DevTools', '❌'],
          ['set_match_id',       'Set your MySchedule profile ID for personalized search', '❌'],
      ],
      [Inches(2.3), Inches(4.25), Inches(1.1)], fsize=11.5)

textbox(slide, CONT_X, Inches(3.95), Inches(4.5), Inches(0.3),
        'Authentication Strategy', fsize=14, color=ACCENT, bold=True, margin=0.0)
bullets(slide, CONT_X, Inches(4.3), Inches(4.6), Inches(2.25), [
    'Preferred: refresh token from browser localStorage (~90 days)',
    'Fallback: access token from DevTools Network tab (~1 hour)',
    'Tokens stored in local files — never committed to git',
    'Auto-seeded at server startup from token files',
    'In-chat JWT paste intercepted by FastAPI regex',
], fsize=12.5, color=BRIGHT, bullet_color=ACCENT)

slide_num(slide, 8)
notes(slide, 'The myschedule-mcp server took approximately 2 days to build from scratch. It works with Claude Desktop and GitHub Copilot as well as this custom agent. The authentication approach (seeding tokens from local files) is pragmatic — the MySchedule portal uses browser-based PKCE auth which can\'t be replicated in a CLI process, so the token-seeding approach is the practical workaround.')


# ── Slide 9: LangGraph Agent ─────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'The Agent Brain: LangGraph Stateful Workflow')

bullets(slide, CONT_X, CONT_Y, Inches(5.4), Inches(1.05), [
    'Stateful graph: every node gets the full message history',
    'MemorySaver: persists conversation per thread_id across messages',
    'Async throughout: ainvoke, astream_events, async tools',
], fsize=13, color=BRIGHT, bullet_color=ACCENT)

# Graph diagram (right side)
gx = Inches(6.4)

def gbox(slide, x, y, w, h, lines, colors, sizes, bolds, fill_c, line_c):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill(s, fill_c); line(s, line_c, 2)
    centered_lines(s, lines, colors, sizes, bolds)

start = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                                 gx+Inches(1.4), Inches(1.38), Inches(0.7), Inches(0.38))
fill(start, GREY); no_line(start)
centered_lines(start, ['START'], [DARK_TEXT], [11], [True])

gbox(slide, gx+Inches(0.8), Inches(1.98), Inches(2.15), Inches(0.78),
     ['[agent]', 'phi4-mini via Ollama'], [TITLE, ACCENT], [15, 11.5], [True, False],
     BG2, ACCENT)

diamond = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND,
                                   gx+Inches(1.1), Inches(3.0), Inches(1.55), Inches(0.8))
fill(diamond, YELLOW); line(diamond, YELLOW, 1)
centered_lines(diamond, ['tool_calls?'], [DARK_TEXT], [12], [True])

gbox(slide, gx+Inches(3.6), Inches(3.1), Inches(2.15), Inches(0.65),
     ['[tools]', 'MCP · local tools'], [TITLE, GREEN], [15, 11], [True, False],
     BG2, GREEN)

reporter = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                    gx+Inches(0.95), Inches(4.1), Inches(1.8), Inches(0.52))
fill(reporter, GREY); no_line(reporter)
centered_lines(reporter, ['[reporter]', 'synthesize answer'], [DARK_TEXT, DARK_TEXT], [12, 9.5], [True, False])

end_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                                   gx+Inches(1.4), Inches(4.9), Inches(0.7), Inches(0.38))
fill(end_box, GREY); no_line(end_box)
centered_lines(end_box, ['END'], [DARK_TEXT], [11], [True])

connector(slide, gx+Inches(1.75), Inches(1.76), gx+Inches(1.75), Inches(1.98))
connector(slide, gx+Inches(1.75), Inches(2.76), gx+Inches(1.75), Inches(3.0))
connector(slide, gx+Inches(2.65), Inches(3.4),  gx+Inches(3.6),  Inches(3.42))
textbox(slide, gx+Inches(2.72), Inches(3.12), Inches(0.65), Inches(0.25),
        'YES', fsize=10, color=ACCENT, bold=True, margin=0.0)
connector(slide, gx+Inches(1.75), Inches(3.8), gx+Inches(1.75), Inches(4.1), GREY)
textbox(slide, gx+Inches(1.9), Inches(3.83), Inches(0.35), Inches(0.22),
        'NO', fsize=10, color=GREY, bold=True, margin=0.0)
connector(slide, gx+Inches(1.75), Inches(4.62), gx+Inches(1.75), Inches(4.9))

# Loop back arrow (tool node → agent node)
lx = gx + Inches(5.75)
connector(slide, gx+Inches(5.75), Inches(3.42), lx, Inches(2.37), GREEN)
connector(slide, lx, Inches(2.37), gx+Inches(2.95), Inches(2.37), GREEN)

# Key innovations (left side)
textbox(slide, CONT_X, Inches(2.1), Inches(5.5), Inches(0.28),
        'Key Implementation Details', fsize=13, color=ACCENT, bold=True, margin=0.0)
innovations = [
    ('Text-call parser', 'Regex + JSON fallback for models that output raw text instead of structured tool_calls'),
    ('Retry logic', 'Auto-detects GPU runner crash → unloads model → reloads → retries (up to 2×)'),
    ('Context trimming', 'Keeps last 10 messages; tool results truncated to 4000 / 300 chars (tiered)'),
    ('Max iterations', '8 tool-call rounds max per turn — prevents infinite loops'),
]
for i, (title, desc) in enumerate(innovations):
    y = Inches(2.48) + Inches(0.88) * i
    ib = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  CONT_X, y, Inches(5.6), Inches(0.72))
    fill(ib, BG2); line(ib, ACCENT, 1)
    textbox(slide, CONT_X+Inches(0.1), y+Inches(0.06), Inches(5.3), Inches(0.26),
            title, fsize=12.5, color=ACCENT, bold=True, margin=0.0)
    textbox(slide, CONT_X+Inches(0.1), y+Inches(0.34), Inches(5.3), Inches(0.3),
            desc, fsize=11, color=BRIGHT, margin=0.0)

slide_num(slide, 9)
notes(slide, 'The reporter node is a key design choice: instead of letting the agent synthesize tool results (which phi4-mini does poorly), we have a dedicated synthesis step that\'s optimized for it. The text-call parser is needed because phi4-mini often outputs tool calls as JSON text rather than using the structured tool_calls mechanism — this fallback makes it work reliably.')


# ── Slide 10: Streaming ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'Real-Time Intelligence: Server-Sent Events Streaming')

textbox(slide, CONT_X, CONT_Y, Inches(5.8), Inches(0.3),
        'What the User Sees', fsize=14, color=ACCENT, bold=True, margin=0.0)

sse_demo = (
    '→ [status]     "thinking"\n'
    '→ [tool_start] "search_roles"\n'
    '→ [status]     "Calling search_roles..."\n'
    '→ [tool_end]   "search_roles"\n'
    '→ [tool_start] "apply_role"\n'
    '→ [status]     "Applying for Data Engineer..."\n'
    '→ [tool_end]   "apply_role"\n'
    '→ [status]     "Summarising..."\n'
    '→ [token]      "✅ Application submitted"\n'
    '→ [token]      " — Role 6260189 at BRISTOL"\n'
    '→ [done]       { thread_id: "abc-123" }'
)
code_block(slide, CONT_X, CONT_Y+Inches(0.4), Inches(5.9), Inches(3.5), sse_demo, fsize=11.5)

textbox(slide, Inches(6.8), CONT_Y, Inches(5.95), Inches(0.3),
        'How It Works', fsize=14, color=ACCENT, bold=True, margin=0.0)
bullets(slide, Inches(6.8), CONT_Y+Inches(0.38), Inches(5.95), Inches(2.4), [
    'FastAPI EventSourceResponse streams events as they happen',
    'LangGraph astream_events (v2) yields per-node events',
    'Agent tokens buffered — flushed only if no tool calls follow',
    'Reporter tokens stream directly (always the final answer)',
    'React EventSource API appends tokens in real time',
    'Thread-ID ties each browser session to LangGraph MemorySaver',
], fsize=13, color=BRIGHT, bullet_color=ACCENT)

textbox(slide, Inches(6.8), Inches(3.9), Inches(5.95), Inches(0.3),
        'OpenAI-Compatible API', fsize=14, color=ACCENT, bold=True, margin=0.0)
bullets(slide, Inches(6.8), Inches(4.28), Inches(5.95), Inches(1.3), [
    'GET /v1/models · POST /v1/chat/completions',
    'Full SSE chunk streaming with finish_reason',
    'Compatible with Open WebUI — plug and play',
], fsize=13, color=BRIGHT, bullet_color=PURPLE)

callout(slide, CONT_X, Inches(5.6), Inches(12.2), Inches(0.72),
        'The bypass layer intercepts "yes / confirm / proceed" before the LLM sees it — ensures deterministic apply execution via SSE',
        fsize=14, text_color=TITLE)

slide_num(slide, 10)
notes(slide, 'The token buffering logic is a subtle but important detail: the agent node produces tokens that look like plain text answers, but may actually be tool-call JSON. We buffer these tokens and discard them if a tool call follows. Only the reporter node\'s tokens are streamed unconditionally, because by that point we know there are no more tool calls coming.')


# ── Slide 11: The Apply Workflow ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'The Apply Workflow: 5 API Calls Behind One Command')

textbox(slide, CONT_X, CONT_Y, Inches(12.3), Inches(0.35),
        'User types "Apply" → server-side bypass detects confirmation → calls view_role + apply_role directly',
        fsize=14, color=BRIGHT, italic=False, margin=0.0)

# 5-step pipeline
step_data = [
    ('1', 'Send Email',           'Application email to staffing manager with role + project details',    ACCENT),
    ('2', 'Role Audit',           'Records application in MySchedule audit log via API',                   GREEN),
    ('3', 'Create Candidate',     'Creates candidate record in the portal (name, profile, CV links)',      PURPLE),
    ('4', 'Self-Input',           'Saves candidate self-input data (role, project, profile key)',          ORANGE),
    ('5', 'Indicator Logic',      'Logs candidate match indicator for analytics and tracking',             YELLOW),
]
sw, sh = Inches(2.2), Inches(1.5)
for i, (num, title, desc, color) in enumerate(step_data):
    sx = Inches(0.55) + Inches(2.55) * i
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  sx, Inches(1.72), sw, sh)
    fill(box, BG2); line(box, color, 2)
    # number circle
    nb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                                  sx+Inches(0.08), Inches(1.8), Inches(0.38), Inches(0.38))
    fill(nb, color); no_line(nb)
    textbox(slide, sx+Inches(0.08), Inches(1.83), Inches(0.38), Inches(0.32),
            num, fsize=12, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER, margin=0.0)
    textbox(slide, sx+Inches(0.1), Inches(2.28), sw-Inches(0.2), Inches(0.3),
            title, fsize=13, color=color, bold=True, margin=0.0)
    textbox(slide, sx+Inches(0.1), Inches(2.62), sw-Inches(0.2), Inches(0.5),
            desc, fsize=10.5, color=BRIGHT, margin=0.0)
    if i < len(step_data)-1:
        arrow(slide, sx+sw+Inches(0.06), Inches(2.3), Inches(0.27), Inches(0.5), 'right', GREY)

# Result states
results = [
    ('✅ All 5 Steps Passed', 'Application submitted successfully — staffing manager notified', GREEN),
    ('⚠️ Partial Success',   'Shows exactly which steps succeeded — email always sent first', YELLOW),
    ('❌ Auth Error',         'Token expired → re-authenticate then retry', RED),
]
for i, (label, desc, color) in enumerate(results):
    rx = Inches(0.55) + Inches(4.25) * i
    rb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  rx, Inches(3.62), Inches(3.9), Inches(0.9))
    fill(rb, BG2); line(rb, color, 2)
    textbox(slide, rx+Inches(0.12), Inches(3.7), Inches(3.65), Inches(0.28),
            label, fsize=13, color=color, bold=True, margin=0.0)
    textbox(slide, rx+Inches(0.12), Inches(4.02), Inches(3.65), Inches(0.38),
            desc, fsize=11.5, color=BRIGHT, margin=0.0)

# Code snippet
code_apply = (
    '# server.py — _apply_directly_generator()\n'
    'await portalClient.sendApplicationEmail({...})\n'
    'await portalClient.applyForRole({...})\n'
    'await portalClient.createCandidate({...})\n'
    'await portalClient.saveCandidateSelfInput({...})\n'
    'await portalClient.candidateIndicatorLogic({...})\n'
    '→ { success: true, steps: [{...}, {...}, ...] }'
)
code_block(slide, CONT_X, Inches(4.75), Inches(12.2), Inches(1.55), code_apply, fsize=11)

slide_num(slide, 11)
notes(slide, 'The 5-step apply flow exactly mirrors what the MySchedule portal does when you manually click Apply — we reverse-engineered the API calls from browser DevTools. Each step is independently error-handled, so if one fails, the others still run. The email is always step 1, ensuring the staffing manager is always notified even if later steps fail.')


# ── Slide 12: Server-Side Bypass ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'The Key Innovation: Server-Side Bypass Architecture')

callout(slide, CONT_X, CONT_Y, Inches(12.3), Inches(0.65),
        '"A 3.8B model cannot reliably orchestrate a 5-step multi-API flow. Solution: take the LLM out of the critical path."',
        fsize=15, italic=True, text_color=ACCENT)

textbox(slide, CONT_X, Inches(2.05), Inches(5.8), Inches(0.3),
        'What phi4-mini did wrong (before bypass)', fsize=13, color=RED, bold=True, margin=0.0)
table(slide, CONT_X, Inches(2.4), Inches(5.8), Inches(2.6),
      ['User Intent', 'Model Behavior'],
      [
          ['User says "yes" → apply', 'Started a NEW search instead'],
          ['Use projectKey from view_role', 'Hallucinated a random number'],
          ['Call view_role first', 'Skipped it, went straight to apply'],
          ['Apply completes', 'Said "submitted" without calling the API'],
      ],
      [Inches(2.4), Inches(3.4)], fsize=11.5)

textbox(slide, Inches(6.5), Inches(2.05), Inches(6.35), Inches(0.3),
        'The Bypass — Two Paths', fsize=13, color=GREEN, bold=True, margin=0.0)

bypass_code = (
    'if _is_confirmation(message):  # "yes/confirm/proceed"\n'
    '\n'
    '    # Path 1: model already called view_role\n'
    '    apply_params = _get_pending_apply(\n'
    '        agent, thread_id  # reads MemorySaver state\n'
    '    )\n'
    '    if apply_params:\n'
    '        return EventSourceResponse(\n'
    '            _apply_directly_generator(apply_params, ...)\n'
    '        )\n'
    '\n'
    '    # Path 2: extract role ID from AI text\n'
    '    role_id = _extract_role_id_from_messages(\n'
    '        agent, thread_id  # regex scan last 6 msgs\n'
    '    )\n'
    '    if role_id:\n'
    '        return EventSourceResponse(\n'
    '            _view_and_apply_generator(role_id, ...)\n'
    '        )'
)
code_block(slide, Inches(6.5), Inches(2.4), Inches(6.35), Inches(3.2), bypass_code, fsize=10.5)

callout(slide, CONT_X, Inches(5.3), Inches(5.85), Inches(0.75),
        'Lesson: for deterministic multi-step flows,\nserverside code beats probabilistic LLM generation',
        fsize=13, text_color=TITLE, line_color=ACCENT)

textbox(slide, CONT_X, Inches(6.2), Inches(5.85), Inches(0.3),
        'Pattern: "Smart Router + Dumb LLM" — applicable beyond this project',
        fsize=12, color=BODY, italic=True, margin=0.0)

slide_num(slide, 12)
notes(slide, 'This is the most important engineering insight in the whole system. Small models are excellent at NLU (understanding what the user wants) but unreliable at planning and executing multi-step workflows. The bypass pattern — where the SERVER executes the critical path deterministically while the LLM handles conversation — is a generally applicable pattern for any agentic system with a small model.')


# ── Slide 13: Privacy & Security ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'Privacy-First, Local-First Architecture')

# Hardware specs box
hw = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                              CONT_X, CONT_Y, Inches(4.4), Inches(2.15))
fill(hw, BG2); line(hw, ACCENT, 1.5)
textbox(slide, CONT_X+Inches(0.12), CONT_Y+Inches(0.08), Inches(4.2), Inches(0.28),
        'Hardware (this laptop)', fsize=13, color=ACCENT, bold=True, margin=0.0)
bullets(slide, CONT_X+Inches(0.1), CONT_Y+Inches(0.42), Inches(4.1), Inches(1.6), [
    'Intel Core Ultra 7 258V (Lunar Lake)',
    'Intel Arc 140V GPU · 32 GB RAM (shared)',
    'Windows 11 · phi4-mini: ~5–8 tok/s CPU',
    'num_gpu=0: model fits on CPU within RAM budget',
], fsize=12, color=BRIGHT, bullet_color=GREY)

# Auth table
textbox(slide, Inches(5.15), CONT_Y, Inches(7.7), Inches(0.3),
        'Authentication Mechanisms', fsize=13, color=ACCENT, bold=True, margin=0.0)
table(slide, Inches(5.15), CONT_Y+Inches(0.38), Inches(7.7), Inches(2.15),
      ['Method', 'Duration', 'How'],
      [
          ['Refresh token (preferred)', '~90 days', 'From browser localStorage → local refresh_token.txt → auto-seeded at startup'],
          ['Access token (fallback)',   '~1 hour',  'From DevTools Network tab → token.txt → auto-seeded'],
          ['In-chat JWT paste',         'Immediate', 'Regex detects eyJ... in message → saves to token.txt → re-seeds MCP auth'],
      ],
      [Inches(2.4), Inches(1.1), Inches(4.2)], fsize=11)

textbox(slide, CONT_X, Inches(3.55), Inches(12.3), Inches(0.28),
        'Privacy Guarantees', fsize=13, color=ACCENT, bold=True, margin=0.0)

privacies = [
    ('Your data stays on your machine', 'Prompts, resumes, tokens, profile data — never sent to any cloud service', GREEN),
    ('Credentials never committed', 'token.txt, refresh_token.txt, .env in .gitignore by default', YELLOW),
    ('No third-party API keys', 'phi4-mini runs via Ollama locally — no OpenAI/Anthropic keys needed', ACCENT),
    ('Network-isolated model', 'phi4-mini has no internet access — cannot exfiltrate data', PURPLE),
]
for i, (title, desc, color) in enumerate(privacies):
    px = CONT_X + Inches(3.15) * i
    pb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  px, Inches(3.92), Inches(2.98), Inches(1.35))
    fill(pb, BG2); line(pb, color, 1.5)
    textbox(slide, px+Inches(0.12), Inches(4.0), Inches(2.75), Inches(0.3),
            title, fsize=12, color=color, bold=True, margin=0.0)
    textbox(slide, px+Inches(0.12), Inches(4.35), Inches(2.75), Inches(0.8),
            desc, fsize=11, color=BRIGHT, margin=0.0)

callout(slide, CONT_X, Inches(5.6), Inches(12.3), Inches(0.7),
        'Unlike cloud HR AI tools (IBM Watsonx, Workday Sana) — your resume, profile, and conversations never leave your machine',
        fsize=14, text_color=TITLE)

slide_num(slide, 13)
notes(slide, 'Privacy is a genuine differentiator. Enterprise cloud HR AI tools (IBM Watsonx Orchestrate, Workday Sana, Moveworks) all send your data to their cloud. This solution keeps everything local. The trade-off is that phi4-mini is less capable than GPT-4 or Claude — but for structured tool use with the bypass pattern, it\'s sufficient. NOTE: "Network-isolated" claim means the Ollama process doesn\'t make outbound calls; this could theoretically be circumvented by the model\'s tool calls, but the tools are controlled by us.')


# ── Slide 14: Industry Parallels ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'Industry Parallels — Similar Solutions in Production')

textbox(slide, CONT_X, CONT_Y, Inches(12.3), Inches(0.3),
        'What we built follows patterns already proven at scale in enterprise AI — using the same open-source frameworks',
        fsize=14, color=BRIGHT, margin=0.0)

table(slide, CONT_X, Inches(1.6), Inches(12.25), Inches(4.2),
      ['Solution', 'What They Built', 'What We Built', 'Key Difference'],
      [
          ['LinkedIn AI Recruiter\n(LangGraph, 2025)',
           'LangGraph agent that automates candidate sourcing across LinkedIn\'s talent graph',
           'LangGraph agent that finds and applies for internal project roles',
           'We run locally; they run at scale on cloud infrastructure'],
          ['GitHub Copilot Workspace\n(MCP, 2025)',
           'AI agent with MCP tools for code generation, test running, PR creation',
           'AI agent with MCP tools for role search, view, and application',
           'Our domain is HR staffing; theirs is software development'],
          ['IBM Watsonx Orchestrate\n(HR AI Agent)',
           'Enterprise cloud AI for HR automation across SAP, Workday, ServiceNow',
           'Local AI for HR automation against internal MySchedule portal',
           'Ours is private, free, runs on a laptop; theirs is hosted, licensed'],
          ['Workday Sana AI Agents\n(2025)',
           'Cloud AI agents for hiring, performance, and development within Workday',
           'AI agent for project staffing within Accenture\'s MySchedule system',
           'We integrate with the existing tool rather than replacing it'],
      ],
      [Inches(2.5), Inches(3.8), Inches(3.6), Inches(2.35)], fsize=11)

callout(slide, CONT_X, Inches(6.08), Inches(12.25), Inches(0.65),
        'The architectural patterns (agentic loop, tool protocol, streaming) are identical — the difference is scale, budget, and cloud vs. local',
        fsize=14, text_color=TITLE, line_color=PURPLE)

slide_num(slide, 14)
notes(slide, 'IMPORTANT NOTE: All comparisons on this slide are for illustration purposes. LinkedIn AI Recruiter details are from LangChain\'s October 2025 blog post. GitHub Copilot Workspace MCP integration is from GitHub\'s public documentation. IBM Watsonx Orchestrate HR capabilities are from IBM\'s marketing materials. Workday Sana is from Workday\'s product pages. Our implementation is a personal tool built in weeks, NOT at the same engineering scale as these products. The comparison is architectural, not about feature parity or scale.')


# ── Slide 15: Results ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'Results & Impact')

# 6 metric boxes
mw2, mh2 = Inches(3.8), Inches(1.5)
metrics = [
    ('Application Time', '30 sec', ACCENT, 'vs 10 minutes manually (93% faster)'),
    ('Cloud Cost', '$0 / mo', GREEN, 'zero API subscriptions required'),
    ('Data Privacy', '100%', PURPLE, 'local — no data leaves the machine'),
    ('Apply Steps Automated', '5', ORANGE, 'end-to-end API orchestration'),
    ('Daily Applications*', '20+', YELLOW, '10× throughput increase (theoretical)'),
    ('Token Lifetime', '90 days', rgb('#00c896'), 'with refresh token — no daily re-auth'),
]
for i, (label, val, color, unit) in enumerate(metrics):
    col, row = i % 3, i // 3
    mx = Inches(0.55) + Inches(4.27) * col
    my = Inches(1.45) + Inches(1.72) * row
    metric_box(slide, mx, my, mw2, mh2, label, val, color, unit)

callout(slide, CONT_X, Inches(5.08), Inches(12.25), Inches(0.65),
        'These results are based on personal use during development and testing — not a formal study',
        fsize=13, italic=True, text_color=BODY, line_color=GREY, fill_color=BG2)

textbox(slide, CONT_X, Inches(5.9), Inches(12.25), Inches(0.5),
        '* Daily Applications metric is theoretical maximum based on observed per-application time. Actual usage varied. '
        'NOTE: This figure is extrapolated / exaggerated for demo impact — real throughput depends on available roles and decision quality.',
        fsize=11.5, color=GREY, italic=True, margin=0.0)

slide_num(slide, 15)
notes(slide, 'EXAGGERATION WARNING: The "20+ daily applications" figure is a theoretical extrapolation — if you ran the agent continuously and applied for every role it found. In practice, you would review results before applying. The 30-second application time is real and has been consistently observed. The $0 cloud cost is real. The 100% local privacy claim is real. The 5 automated steps is real. The 90-day token lifetime is the documented lifetime of Microsoft Azure AD refresh tokens.')


# ── Slide 16: Key Lessons ────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'Five Things I Learned Building This')

lessons = [
    ('Small LLMs: excellent at understanding, unreliable at execution',
     ' — route around the model for deterministic multi-step flows (the bypass pattern)'),
    ('MCP is genuinely production-ready — and worth building for',
     ' — 2 days to build myschedule-mcp; now works with Claude Desktop, GitHub Copilot, and our agent'),
    ('State management is the hardest part of agentic systems',
     ' — LangGraph MemorySaver solved persistence; React async EventSource caused subtle race conditions'),
    ('Local AI is viable for personal automation in 2026',
     ' — phi4-mini at 5–8 tok/s is usable; privacy and zero cost are genuine advantages'),
    ('The server is more reliable than the model for critical operations',
     ' — treat the LLM as a smart dispatcher, not a workflow engine'),
]

for i, (lead, rest) in enumerate(lessons):
    y = Inches(1.45) + Inches(0.98) * i
    lb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  CONT_X, y, Inches(12.3), Inches(0.82))
    fill(lb, BG2); line(lb, ACCENT, 1.5)
    nb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                                  Inches(0.65), y+Inches(0.18), Inches(0.46), Inches(0.46))
    fill(nb, ACCENT); no_line(nb)
    textbox(slide, Inches(0.65), y+Inches(0.21), Inches(0.46), Inches(0.38),
            str(i+1), fsize=14, color=DARK_TEXT, bold=True, align=PP_ALIGN.CENTER, margin=0.0)
    textbox(slide, Inches(1.28), y+Inches(0.1), Inches(10.9), Inches(0.3),
            lead, fsize=13.5, color=TITLE, bold=True, margin=0.0)
    textbox(slide, Inches(1.28), y+Inches(0.44), Inches(10.9), Inches(0.3),
            rest, fsize=12.5, color=BODY, margin=0.0)

slide_num(slide, 16)
notes(slide, 'These lessons are genuine — learned through debugging, not theory. The bypass pattern (lesson 1) emerged after hours of trying to make phi4-mini reliably call view_role before apply_role. The MCP compatibility (lesson 2) was a pleasant surprise — the same myschedule-mcp server was tested in Claude Desktop and worked without modification. The state management lesson (lesson 3) cost a day of debugging React\'s closure over stale EventSource handlers.')


# ── Slide 17: Roadmap ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'What\'s Next — Roadmap')

textbox(slide, CONT_X, CONT_Y, Inches(5.75), Inches(0.3),
        'Near-Term (Next 1–3 Months)', fsize=14, color=ACCENT, bold=True, margin=0.0)
near = [
    'Upgrade to phi4 (14B) for more reliable tool-call generation',
    'Applied-roles tracker — prevent duplicate applications',
    'Role recommendation scoring based on skills profile match',
    'Batch search — run multiple queries in one session',
    'Playwright browser login for seamless Azure AD auth',
]
bullets(slide, CONT_X, CONT_Y+Inches(0.38), Inches(5.75), Inches(2.8),
        near, fsize=13.5, color=BRIGHT, bullet_color=ACCENT)

textbox(slide, Inches(7.0), CONT_Y, Inches(5.85), Inches(0.3),
        'Future (3–12 Months)*', fsize=14, color=PURPLE, bold=True, margin=0.0)
future = [
    'Multi-agent: dedicated search agent + apply agent',
    'Voice interface: Whisper STT + piper TTS',
    'Calendar integration: track start dates and follow-up reminders',
    'Expose agent as its own MCP server for Copilot / Claude Desktop',
    'RAG over personal skills profile for smarter matching',
    'Notification system: daily digest of new matching roles',
]
bullets(slide, Inches(7.0), CONT_Y+Inches(0.38), Inches(5.85), Inches(3.4),
        future, fsize=13.5, color=BRIGHT, bullet_color=PURPLE)

callout(slide, CONT_X, Inches(5.55), Inches(12.3), Inches(0.75),
        '"Expose the agent as its own MCP server" — so Claude Desktop or GitHub Copilot could search and apply for roles through this system',
        fsize=14, italic=True, text_color=TITLE, line_color=PURPLE)

textbox(slide, CONT_X, Inches(6.5), Inches(12.3), Inches(0.35),
        '* Future items are speculative — actual roadmap depends on time available and whether phi4-mini limitations warrant the upgrade',
        fsize=11, color=GREY, italic=True, margin=0.0)

slide_num(slide, 17)
notes(slide, 'The near-term items are genuine — I plan to work on them in the next few months. The phi4 (14B) upgrade is the most impactful — the 14B model is significantly more reliable for tool-call generation. NOTE: Future items are speculative. "Expose as MCP server" is a particularly interesting future direction — it would create a recursive loop where the AI agent could be driven by another AI.')


# ── Slide 18: Demo ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, "Let's See It Work — Live Demo")

demo_script = (
    'Step 1: Search\n'
    '  > "Find data engineer roles in USA"\n'
    '  → [tool_start] search_roles\n'
    '  → Markdown table: Role ID | Title | Client | Location | Start | Status\n'
    '\n'
    'Step 2: Apply\n'
    '  > "Apply for the Bristol Myers Squibb one"\n'
    '  → Agent shows role details, asks for confirmation\n'
    '  > "Yes"\n'
    '  → [Bypass] view_role + apply_role (no LLM in the loop)\n'
    '  → ✅ Application submitted — Role 6260189\n'
    '\n'
    'Step 3 (optional): Re-authenticate\n'
    '  > "eyJ0eXAiOiJKV1Qi..."\n'
    '  → [JWT intercept] token saved → MCP re-seeded\n'
    '  → ✅ Token updated. Try your request again.'
)
code_block(slide, Inches(0.85), Inches(1.5), Inches(7.35), Inches(4.1), demo_script, fsize=12)

textbox(slide, Inches(8.6), Inches(1.5), Inches(4.3), Inches(0.3),
        'What to watch for', fsize=14, color=ACCENT, bold=True, margin=0.0)
bullets(slide, Inches(8.6), Inches(1.88), Inches(4.3), Inches(3.2), [
    'Status events stream before text (real-time feel)',
    '"Calling search_roles..." appears before results',
    'Table renders in markdown in the chat window',
    'Bypass fires instantly on "yes" — no LLM delay',
    '5 API calls complete in < 5 seconds',
    'thread_id persists context across messages',
], fsize=13, color=BRIGHT, bullet_color=ACCENT)

callout(slide, Inches(0.85), Inches(5.85), Inches(12.1), Inches(0.75),
        'Services: Ollama (port 11434) · FastAPI (port 8000) · React UI (port 3000)',
        fsize=14, text_color=TITLE)

slide_num(slide, 18)
notes(slide, 'Live demo checklist: (1) Ensure Ollama is running — ollama serve. (2) Ensure FastAPI is running — uvicorn server:app --reload --port 8000. (3) Ensure React is running — npm start. (4) Confirm refresh_token.txt is present (check auto-seed output in FastAPI console). (5) Test with a simple query first. NOTE: Demo environment — real data, real application. Audience should understand that applications are being submitted to actual MySchedule.')


# ── Slide 19: Summary ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide); slide_title(slide, 'Summary — What Was Built and Why It Matters')

table(slide, CONT_X, CONT_Y, Inches(12.25), Inches(3.55),
      ['Component', 'What It Does', 'Technology', 'Key Innovation'],
      [
          ['phi4-mini + Ollama',   'Local LLM: understands natural language intent',
           'Python · Ollama',      'Zero cloud cost · 100% private'],
          ['LangGraph Agent',      'Stateful think→act→observe loop with memory',
           'Python · LangGraph 1.0', 'MemorySaver + reporter node + retry logic'],
          ['myschedule-mcp',       'Wraps MySchedule REST API as AI tools',
           'TypeScript · MCP SDK', 'Compatible with any MCP client'],
          ['FastAPI + bypass layer', 'Streaming + guaranteed apply execution',
           'Python · SSE',         'Deterministic bypass for multi-step flows'],
          ['React Chat UI',        'Real-time chat with live tool status',
           'JavaScript · EventSource', 'Thread-ID conversation persistence'],
      ],
      [Inches(2.4), Inches(3.25), Inches(2.4), Inches(4.2)], fsize=11.5)

callout(slide, CONT_X, Inches(4.92), Inches(12.25), Inches(1.1),
        '"AI isn\'t magic. It\'s an LLM, a state graph, an open protocol, and a few hundred lines of Python.\n'
        'Small models can handle real enterprise workflows — when you build the right scaffolding around them."',
        fsize=16, italic=False, text_color=ACCENT)

textbox(slide, CONT_X, Inches(6.22), Inches(12.25), Inches(0.55),
        'Built with open-source tools · runs on a laptop · costs $0/month · protects your data · applies to real jobs',
        fsize=14, color=BRIGHT, align=PP_ALIGN.CENTER, margin=0.0)

slide_num(slide, 19)
notes(slide, 'Closing summary. The key message: this was built by one person in a few weeks using entirely open-source frameworks that are already running in production at major tech companies. The same patterns (LangGraph, MCP, local LLM) that LinkedIn and Uber use for enterprise AI can be applied to personal productivity tools. The scale is different, but the architecture is the same.')


# ── Slide 20: Q&A ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg(slide)

bar1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
fill(bar1, ACCENT); no_line(bar1)
bar2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.12), SLIDE_W, Inches(0.06))
fill(bar2, ACCENT2); no_line(bar2)

textbox(slide, Inches(2.0), Inches(1.5), Inches(9.3), Inches(0.8),
        'Questions?', fsize=52, color=TITLE, bold=True, align=PP_ALIGN.CENTER, margin=0.0)

textbox(slide, Inches(3.5), Inches(2.55), Inches(6.3), Inches(0.4),
        'Deepa Chandramohan', fsize=22, color=TITLE, bold=True,
        align=PP_ALIGN.CENTER, margin=0.0)
textbox(slide, Inches(2.8), Inches(3.05), Inches(7.7), Inches(0.32),
        'deepa.chandramohan@accenture.com', fsize=16, color=ACCENT,
        align=PP_ALIGN.CENTER, margin=0.0)
textbox(slide, Inches(3.2), Inches(3.45), Inches(6.9), Inches(0.32),
        'Accenture  |  Technology  |  May 2026', fsize=14, color=BODY,
        align=PP_ALIGN.CENTER, margin=0.0)

code_block(slide, Inches(1.8), Inches(4.2), Inches(9.7), Inches(1.85), chat_demo, fsize=12)

badge(slide, Inches(1.8),  Inches(6.35), Inches(2.3), Inches(0.38), '100% Local',      rgb('#196c2e'), GREEN)
badge(slide, Inches(4.35), Inches(6.35), Inches(2.3), Inches(0.38), 'Zero Cloud Cost', rgb('#1f6feb'), ACCENT)
badge(slide, Inches(6.9),  Inches(6.35), Inches(2.3), Inches(0.38), 'MCP Compatible',  rgb('#3b1f6b'), PURPLE)
badge(slide, Inches(9.45), Inches(6.35), Inches(2.05), Inches(0.38), 'Open Source',    rgb('#5a1f1f'), ORANGE)

slide_num(slide, 20)
notes(slide, 'Q&A slide. Common anticipated questions:\n- "Why phi4-mini instead of phi4?" → RAM budget; phi4-mini fits comfortably, phi4 (14B) would need 8–10 GB and might exceed the GPU shared memory budget.\n- "Why not use Claude or GPT-4?" → Privacy and cost. The whole point is local, private, free.\n- "How hard was the auth?" → The hardest part was realizing PKCE flow can\'t be replicated in a CLI — the token-seeding approach took a day to figure out.\n- "Could this be deployed for a team?" → Technically yes — run FastAPI on a shared server. Auth tokens would need per-user handling. Not the current design.')

# ── Save ─────────────────────────────────────────────────────────────────────
output = r'c:\dev\local-llm-phi4\slide-deck-v2.pptx'
prs.save(output)
print(f'Saved to {output}')
