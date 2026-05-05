from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.replace('#', '')
    return RGBColor.from_string(hex_value.upper())


BG = rgb('#1e2030')
ACCENT = rgb('#00d4ff')
TITLE = rgb('#ffffff')
BODY = rgb('#e0e0e0')
CALLOUT_BG = rgb('#2a2d45')
CALLOUT_BORDER = ACCENT
CODE_BG = rgb('#0d1117')
CODE_TEXT = rgb('#c9d1d9')
TABLE_HEADER_BG = ACCENT
TABLE_HEADER_TEXT = BG
TABLE_ROW_1 = rgb('#252840')
TABLE_ROW_2 = rgb('#2a2d45')
RED = rgb('#ff6b6b')
YELLOW = rgb('#ffd700')
GREEN = rgb('#4caf50')
GREY = rgb('#9aa4b2')
DARK_TEXT = BG

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
TITLE_X = Inches(0.4)
TITLE_Y = Inches(0.3)
TITLE_W = Inches(12.5)
TITLE_H = Inches(0.8)
CONTENT_X = Inches(0.4)
CONTENT_Y = Inches(1.3)
CONTENT_W = Inches(12.5)
CONTENT_H = Inches(5.7)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def apply_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def apply_line(shape, color, width=1.5):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def style_run(run, *, name='Calibri', size=16, color=BODY, bold=False, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_textbox(slide, left, top, width, height, text='', *, font_name='Calibri', font_size=16,
                color=BODY, bold=False, italic=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                margin=0.05):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    style_run(run, name=font_name, size=font_size, color=color, bold=bold, italic=italic)
    return box


def add_title(slide, title):
    return add_textbox(
        slide, TITLE_X, TITLE_Y, TITLE_W, TITLE_H, title,
        font_name='Calibri', font_size=28, color=TITLE, bold=True,
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, margin=0.0
    )


def add_slide_number(slide, number):
    add_textbox(
        slide, Inches(12.65), Inches(7.05), Inches(0.45), Inches(0.22), str(number),
        font_size=10, color=GREY, align=PP_ALIGN.RIGHT, margin=0.0
    )


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_bullets(slide, left, top, width, height, items, *, font_size=16, color=BODY,
                line_spacing=1.15, bullet_color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.line_spacing = line_spacing
        bullet_run = p.add_run()
        bullet_run.text = '• '
        style_run(bullet_run, size=font_size, color=bullet_color or color)
        text_run = p.add_run()
        text_run.text = item
        style_run(text_run, size=font_size, color=color)
    return box


def add_numbered_paragraphs(slide, left, top, width, height, items, *, font_size=16, color=BODY,
                            lead_bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    for index, item in enumerate(items, start=1):
        p = tf.paragraphs[0] if index == 1 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        num_run = p.add_run()
        num_run.text = f'{index}. '
        style_run(num_run, size=font_size, color=color, bold=True)
        if isinstance(item, tuple):
            lead, desc = item
            lead_run = p.add_run()
            lead_run.text = lead
            style_run(lead_run, size=font_size, color=color, bold=lead_bold)
            desc_run = p.add_run()
            desc_run.text = desc
            style_run(desc_run, size=font_size, color=color)
        else:
            text_run = p.add_run()
            text_run.text = item
            style_run(text_run, size=font_size, color=color)
    return box


def add_callout(slide, left, top, width, height, text, *, font_size=16, italic=False,
                fill_color=CALLOUT_BG, line_color=CALLOUT_BORDER, text_color=TITLE,
                align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    apply_fill(shape, fill_color)
    apply_line(shape, line_color, 2)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    style_run(run, size=font_size, color=text_color, italic=italic)
    return shape


def add_code_block(slide, left, top, width, height, code, *, font_size=12, color=CODE_TEXT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    apply_fill(shape, CODE_BG)
    apply_line(shape, ACCENT, 1.25)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    for index, line in enumerate(code.splitlines()):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        style_run(run, name='Courier New', size=font_size, color=color)
    return shape


def set_cell(cell, text, *, font_size=12, color=BODY, bold=False, fill_color=None, align=PP_ALIGN.LEFT):
    if fill_color is not None:
        apply_fill(cell, fill_color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    style_run(run, size=font_size, color=color, bold=bold)


def add_table(slide, left, top, width, height, headers, rows, col_widths, *, font_size=12):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    for idx, col_w in enumerate(col_widths):
        table.columns[idx].width = col_w
    for idx, header in enumerate(headers):
        set_cell(table.cell(0, idx), header, font_size=13, color=TABLE_HEADER_TEXT, bold=True,
                 fill_color=TABLE_HEADER_BG, align=PP_ALIGN.CENTER)
    for r_index, row in enumerate(rows, start=1):
        fill_color = TABLE_ROW_1 if r_index % 2 == 1 else TABLE_ROW_2
        for c_index, value in enumerate(row):
            align = PP_ALIGN.CENTER if c_index == 0 and len(headers) <= 3 and len(str(value)) < 12 else PP_ALIGN.LEFT
            set_cell(table.cell(r_index, c_index), str(value), font_size=font_size, color=BODY,
                     fill_color=fill_color, align=align)
    return table


def add_labeled_box(slide, left, top, width, height, label, label_color, text, *, text_size=16):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    apply_fill(shape, CALLOUT_BG)
    apply_line(shape, ACCENT, 1.5)
    add_textbox(slide, left + Inches(0.08), top + Inches(0.08), width - Inches(0.16), Inches(0.30), label,
                font_size=16, color=label_color, bold=True)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.30), height - Inches(0.70), text,
                font_size=text_size, color=TITLE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return shape


def add_arrow(slide, left, top, width, height, direction='right', color=ACCENT):
    shape_type = {
        'right': MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        'down': MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
        'left': MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
        'up': MSO_AUTO_SHAPE_TYPE.UP_ARROW,
    }[direction]
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    apply_fill(shape, color)
    shape.line.fill.background()
    return shape


def add_connector_arrow(slide, begin_x, begin_y, end_x, end_y, *, color=ACCENT, width=1.75):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_centered_lines(shape, lines, colors=None, sizes=None, bold=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        style_run(
            run,
            size=(sizes[idx] if sizes else 15),
            color=(colors[idx] if colors else TITLE),
            bold=(bold[idx] if bold else False),
        )


chat_text = (
    '> Find data engineer roles in USA\n'
    '✅ Found 3 roles | Data Engineer | BRISTOL MYERS SQUIBB | New York | Remote\n'
    '> Apply\n'
    '✅ Application submitted — Role 6260189'
)


# Slide 1
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_textbox(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.7), 'From Prompt to Production',
            font_size=36, color=TITLE, bold=True, align=PP_ALIGN.CENTER, margin=0.0)
add_textbox(slide, Inches(1.0), Inches(1.78), Inches(11.3), Inches(0.45),
            'Building a Personal AI Agent with LangGraph, MCP & a Local LLM',
            font_size=20, color=ACCENT, align=PP_ALIGN.CENTER, margin=0.0)
add_textbox(slide, Inches(3.3), Inches(2.30), Inches(6.7), Inches(0.30), 'Deepa Chandramohan | Accenture',
            font_size=14, color=BODY, align=PP_ALIGN.CENTER, margin=0.0)
add_code_block(slide, Inches(2.0), Inches(3.35), Inches(9.3), Inches(1.55), chat_text, font_size=12, color=ACCENT)
add_slide_number(slide, 1)

# Slide 2
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'The Problem with Manual Staffing')
add_bullets(slide, Inches(0.75), Inches(1.65), Inches(11.6), Inches(2.6), [
    'Internal staffing portal (MySchedule) lists hundreds of open project roles',
    'Finding the right role requires repeated searching with different filters',
    'Applying requires: search → view details → confirm → email → audit → candidate record → indicator log',
    '~10 minutes of manual work per application, repeated daily',
], font_size=16, color=TITLE)
add_callout(slide, Inches(1.05), Inches(5.15), Inches(11.1), Inches(0.85),
            'Goal: automate search and apply — without cloud, without sharing credentials',
            font_size=16, text_color=TITLE)
add_notes(slide, 'I wanted to automate this. But I also wanted to learn: can a small, local LLM actually handle real multi-step tool use? So I built the whole thing myself.')
add_slide_number(slide, 2)

# Slide 3
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'What If AI Could Do It?')
add_labeled_box(slide, Inches(1.0), Inches(2.0), Inches(4.4), Inches(1.6), 'Before', RED,
                'You → MySchedule Portal\n(manual, ~10 minutes per role)')
add_arrow(slide, Inches(5.7), Inches(2.45), Inches(1.0), Inches(0.65), 'right', ACCENT)
add_labeled_box(slide, Inches(7.0), Inches(2.0), Inches(5.0), Inches(1.6), 'After', ACCENT,
                'You → AI Agent → MySchedule Portal\n(automatic, ~30 seconds)')
add_bullets(slide, Inches(1.05), Inches(4.35), Inches(10.8), Inches(1.3), [
    'No cloud. No API keys. Runs entirely on your laptop.',
    'Privacy: resume, profile, and auth tokens never leave your machine.',
], font_size=16)
add_slide_number(slide, 3)

# Slide 4
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'The Technology Stack')
add_table(slide, Inches(0.65), Inches(1.55), Inches(12.0), Inches(4.75),
          ['Layer', 'Technology', 'Role'], [
              ['Language Model', 'Microsoft phi4-mini (3.8B)', 'Understands intent, decides which tools to call'],
              ['LLM Runtime', 'Ollama', 'Serves phi4-mini locally on port 11434'],
              ['Agent Orchestration', 'LangGraph', 'Stateful think→act→observe loop'],
              ['AI Framework', 'LangChain', 'Connects LLM, tools, memory'],
              ['Tool Protocol', 'Model Context Protocol (MCP)', 'Standard way to expose external tools to AI'],
              ['Tool Server', 'Node.js/TypeScript MCP server', 'Wraps MySchedule REST API'],
              ['Backend', 'FastAPI (Python)', 'HTTP API + SSE streaming'],
              ['Frontend', 'React', 'Chat interface'],
          ], [Inches(2.25), Inches(3.2), Inches(6.55)], font_size=11.5)
add_textbox(slide, Inches(0.85), Inches(6.55), Inches(11.6), Inches(0.25),
            'Everything except MySchedule itself runs locally — zero cloud spend',
            font_size=10.5, color=GREY, italic=False, align=PP_ALIGN.LEFT, margin=0.0)
add_slide_number(slide, 4)

# Slide 5
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'Model Context Protocol — USB-C for AI Tools')
add_bullets(slide, Inches(0.75), Inches(1.55), Inches(11.6), Inches(2.2), [
    'Open standard released by Anthropic in November 2024',
    'Defines how AI models connect to external tools and data sources',
    'Before MCP: every AI app needed custom integration code per tool',
    'With MCP: build a tool server once — any MCP-compatible client can use it',
], font_size=16)
add_callout(slide, Inches(1.4), Inches(3.6), Inches(10.4), Inches(1.2),
            '"USB-C let any device connect to any charger without custom cables.\nMCP lets any AI model connect to any tool without custom code."',
            font_size=15, italic=True, text_color=TITLE)
add_textbox(slide, Inches(2.0), Inches(5.25), Inches(9.5), Inches(0.4),
            'Used by: Claude Desktop · GitHub Copilot · Cursor · and this agent',
            font_size=16, color=BODY, align=PP_ALIGN.CENTER, margin=0.0)
add_slide_number(slide, 5)

# Slide 6
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'MCP Server: Wrapping MySchedule as AI Tools')
add_textbox(slide, Inches(0.75), Inches(1.65), Inches(4.3), Inches(0.3), 'What it is:',
            font_size=16, color=ACCENT, bold=True, margin=0.0)
add_bullets(slide, Inches(0.75), Inches(2.0), Inches(4.2), Inches(2.3), [
    'Node.js / TypeScript process',
    'Communicates over stdio (subprocess pipes)',
    'Uses @modelcontextprotocol/sdk',
    'Authenticates via Azure AD OAuth2',
], font_size=15)
add_textbox(slide, Inches(5.4), Inches(1.65), Inches(3.0), Inches(0.3), 'Tools exposed:',
            font_size=16, color=ACCENT, bold=True, margin=0.0)
add_table(slide, Inches(5.35), Inches(2.0), Inches(7.0), Inches(3.55),
          ['Tool', 'Description'], [
              ['search_roles', 'Keyword/location search'],
              ['view_role', 'Full details + projectKey'],
              ['apply_role', '5-step application submission'],
              ['seed_token', 'Auth with access token'],
              ['seed_refresh_token', 'Auth with refresh token (90 days)'],
              ['set_match_id', 'Personalised search'],
          ], [Inches(2.35), Inches(4.65)], font_size=12)
add_slide_number(slide, 6)

# Slide 7
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'LangGraph — Stateful Agent Workflows as Graphs')
add_bullets(slide, Inches(0.75), Inches(1.45), Inches(11.8), Inches(1.2), [
    'Library from LangChain for multi-step AI agents',
    'Define logic as a directed graph: nodes = functions, edges = routing',
    'MemorySaver provides conversation persistence across messages (per thread_id)',
], font_size=15)
start = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.85), Inches(2.75), Inches(1.5), Inches(0.45))
apply_fill(start, GREY)
start.line.fill.background()
add_centered_lines(start, ['START'], [BG], [13], [True])
agent_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.15), Inches(3.45), Inches(2.9), Inches(0.95))
apply_fill(agent_box, CALLOUT_BG)
apply_line(agent_box, ACCENT, 2)
add_centered_lines(agent_box, ['[agent node]', 'phi4-mini via Ollama'], [TITLE, ACCENT], [16, 13], [True, False])
diamond = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, Inches(5.55), Inches(4.65), Inches(2.1), Inches(1.1))
apply_fill(diamond, YELLOW)
apply_line(diamond, YELLOW, 1.2)
add_centered_lines(diamond, ['tool_calls?'], [DARK_TEXT], [15], [True])
tools_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(4.75), Inches(2.3), Inches(0.95))
apply_fill(tools_box, CALLOUT_BG)
apply_line(tools_box, GREEN, 2)
add_centered_lines(tools_box, ['[tools node]', 'MCP tools'], [TITLE, GREEN], [16, 13], [True, False])
reporter_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(6.0), Inches(2.4), Inches(0.55))
apply_fill(reporter_box, GREY)
reporter_box.line.fill.background()
add_centered_lines(reporter_box, ['[reporter node]'], [BG], [13], [True])
end_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.95), Inches(6.72), Inches(1.3), Inches(0.38))
apply_fill(end_box, GREY)
end_box.line.fill.background()
add_centered_lines(end_box, ['END'], [BG], [12], [True])
add_connector_arrow(slide, Inches(6.60), Inches(3.20), Inches(6.60), Inches(3.45))
add_connector_arrow(slide, Inches(6.60), Inches(4.40), Inches(6.60), Inches(4.65))
add_connector_arrow(slide, Inches(7.65), Inches(5.20), Inches(9.20), Inches(5.20))
add_textbox(slide, Inches(8.05), Inches(4.75), Inches(0.65), Inches(0.25), 'YES', font_size=11, color=ACCENT, bold=True, margin=0.0)
add_connector_arrow(slide, Inches(6.60), Inches(5.75), Inches(6.60), Inches(6.00), color=GREY)
add_textbox(slide, Inches(6.85), Inches(5.78), Inches(0.55), Inches(0.25), 'NO', font_size=11, color=GREY, bold=True, margin=0.0)
# Manual loop using 3 connectors
c1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.35), Inches(4.75), Inches(10.35), Inches(3.35))
c1.line.color.rgb = GREEN
c1.line.width = Pt(1.75)
c2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(10.35), Inches(3.35), Inches(8.10), Inches(3.35))
c2.line.color.rgb = GREEN
c2.line.width = Pt(1.75)
c3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.10), Inches(3.35), Inches(8.05), Inches(3.95))
c3.line.color.rgb = GREEN
c3.line.width = Pt(1.75)
add_notes(slide, 'The graph handles all the looping. I just define the nodes and routing logic. MemorySaver means the agent remembers everything said in the conversation — across multiple messages.')
add_slide_number(slide, 7)

# Slide 8
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'How Everything Connects')
layers = [
    ('React Chat UI → FastAPI Server (port 8000)', rgb('#1565c0')),
    ('server.py — Bypass Layer (JWT interception, confirmation bypass)', rgb('#6a1b9a')),
    ('LangGraph Agent — phi4-mini via ChatOllama — MemorySaver', rgb('#1b5e20')),
    ('MCP stdio transport — langchain-mcp-adapters', rgb('#e65100')),
    ('myschedule-mcp (Node.js) — search_roles / view_role / apply_role', rgb('#006064')),
    ('MySchedule REST API — Azure AD auth', rgb('#b71c1c')),
]
layer_top = Inches(1.45)
layer_h = Inches(0.64)
for idx, (label, color) in enumerate(layers):
    y = layer_top + Inches(0.83) * idx
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.15), y, Inches(11.0), layer_h)
    apply_fill(rect, color)
    rect.line.fill.background()
    add_centered_lines(rect, [label], [TITLE], [16 if idx < 3 else 15], [True])
    if idx < len(layers) - 1:
        add_arrow(slide, Inches(6.38), y + layer_h + Inches(0.05), Inches(0.55), Inches(0.22), 'down', ACCENT)
add_slide_number(slide, 8)

# Slide 9
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'The Problem with Small Models')
add_table(slide, Inches(0.75), Inches(1.7), Inches(11.8), Inches(3.3),
          ['What We Needed', 'What phi4-mini Did'], [
              ['Search → show table', '✅ Worked reliably'],
              ['Call view_role before applying', '❌ Skipped it, made up projectKey'],
              ['User says "yes" → call apply_role', '❌ Started a NEW search instead'],
              ['Use real projectKey from view_role', '❌ Hallucinated random numbers'],
              ['Apply flow completes', '❌ Said "submitted" without calling the tool'],
          ], [Inches(5.0), Inches(6.8)], font_size=13)
add_callout(slide, Inches(1.1), Inches(5.35), Inches(10.9), Inches(0.95),
            'A 3.8B model is not reliable enough to orchestrate a 5-step apply flow.\nSolution: take the LLM out of the critical path.',
            font_size=16)
add_notes(slide, 'A 3.8B parameter model is not reliable enough to orchestrate a 5-step apply flow. The solution: take the LLM out of the critical path for the apply step.')
add_slide_number(slide, 9)

# Slide 10
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'The Solution: Server-Side Bypass')
add_textbox(slide, Inches(0.75), Inches(1.6), Inches(5.6), Inches(0.35),
            "Trigger: user says 'yes', 'confirm', 'proceed'", font_size=15, color=ACCENT, bold=True, margin=0.0)
add_textbox(slide, Inches(0.75), Inches(2.05), Inches(2.2), Inches(0.3), 'Path 1 — LLM called view_role:',
            font_size=15, color=TITLE, bold=True, margin=0.0)
add_numbered_paragraphs(slide, Inches(0.9), Inches(2.35), Inches(5.5), Inches(1.75), [
    'Read MemorySaver state',
    'Find last view_role ToolMessage',
    'Extract projectKey & projectLocationKey',
    'Call apply_role directly — no LLM',
], font_size=14.5)
add_textbox(slide, Inches(0.75), Inches(4.35), Inches(2.0), Inches(0.3), 'Path 2 — Fallback:',
            font_size=15, color=TITLE, bold=True, margin=0.0)
add_numbered_paragraphs(slide, Inches(0.9), Inches(4.65), Inches(5.5), Inches(1.25), [
    'Regex-scan AI messages for role ID',
    'Call view_role to get real keys',
    'Call apply_role with verified data',
], font_size=14.5)
add_code_block(slide, Inches(6.65), Inches(1.85), Inches(5.1), Inches(3.35), "if _is_confirmation(message):\n    apply_params = _get_pending_apply(\n        agent, thread_id\n    )\n    if apply_params:\n        return EventSourceResponse(\n            _apply_directly_generator(\n                apply_params,\n                thread_id,\n                tool_map\n            )\n        )", font_size=12)
add_notes(slide, "This is the key lesson: don't fight the model's limitations — work around them. For deterministic multi-step flows, the server is more reliable than the LLM.")
add_slide_number(slide, 10)

# Slide 11
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'Inside apply_role — 5 Steps in One Tool Call')
step_titles = [
    'sendEmail\nApplication email to staffing manager',
    'applyToRoleAudit\nRecord in audit log',
    'createCandidate\nCandidate record in portal',
    'saveCandidateSelfInput\nSave self-input data',
    'candidateIndicatorLogic\nLog match indicator',
]
start_x = Inches(0.55)
box_w = Inches(2.12)
for idx, step in enumerate(step_titles):
    x = start_x + Inches(2.45) * idx
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(2.25), box_w, Inches(1.1))
    apply_fill(box, ACCENT)
    box.line.fill.background()
    lines = step.split('\n')
    add_centered_lines(box, [lines[0], lines[1]], [DARK_TEXT, DARK_TEXT], [13, 10.5], [True, False])
    if idx < len(step_titles) - 1:
        add_arrow(slide, x + box_w + Inches(0.08), Inches(2.56), Inches(0.23), Inches(0.35), 'right', ACCENT)
results = [
    ('✅ All passed', 'Application submitted', GREEN),
    ('⚠️ Partial', 'Shows which steps succeeded/failed', YELLOW),
    ('❌ Failed', 'Validation or auth error', RED),
]
for idx, (label, desc, color) in enumerate(results):
    x = Inches(1.0) + Inches(3.95) * idx
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(4.65), Inches(3.4), Inches(0.9))
    apply_fill(box, CALLOUT_BG)
    apply_line(box, color, 2)
    add_centered_lines(box, [label, desc], [color, TITLE], [14, 12], [True, False])
add_slide_number(slide, 11)

# Slide 12
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'Authentication Without Storing Secrets')
add_table(slide, Inches(0.75), Inches(1.55), Inches(11.8), Inches(2.35),
          ['Mechanism', 'Duration', 'How Used'], [
              ['Access token (token.txt)', '~1 hour', 'Pasted from browser DevTools; auto-seeded at server startup'],
              ['Refresh token (refresh_token.txt)', '~90 days', 'From browser localStorage; preferred; survives restarts'],
              ['In-chat token paste', 'Immediate', 'Bypass intercepts JWT → re-seeds without LLM seeing it'],
          ], [Inches(3.1), Inches(1.8), Inches(6.9)], font_size=12)
add_bullets(slide, Inches(0.95), Inches(4.35), Inches(11.0), Inches(1.55), [
    'token.txt and refresh_token.txt are in .gitignore — never committed',
    'All auth is local — tokens never leave the machine',
    'In-chat JWT interception: regex detects eyJ... → save to token.txt → re-seed MCP auth',
], font_size=15)
add_slide_number(slide, 12)

# Slide 13
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'Real-Time Streaming with Server-Sent Events')
add_bullets(slide, Inches(0.75), Inches(1.75), Inches(5.75), Inches(2.7), [
    'FastAPI uses EventSourceResponse to stream agent progress',
    'Events: status (tool in progress), token (text), done (thread_id)',
    'React EventSource appends tokens as they arrive — real-time feel',
    'Each conversation has a thread_id → maps to LangGraph MemorySaver',
], font_size=15)
add_code_block(slide, Inches(6.65), Inches(1.95), Inches(5.1), Inches(2.65),
               '→ [status] "Searching for roles..."\n→ [status] "Calling view_role..."\n→ [status] "Applying for Data Engineer..."\n→ [token]  "✅ Application submitted\n             — Role 6260189..."\n→ [done]   { thread_id: "abc-123" }', font_size=12)
add_slide_number(slide, 13)

# Slide 14
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'Local-First: Everything on Your Laptop')
left_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(5.0), Inches(1.35))
apply_fill(left_box, CALLOUT_BG)
apply_line(left_box, ACCENT, 1.5)
add_textbox(slide, Inches(1.05), Inches(2.0), Inches(4.5), Inches(0.9),
            'Intel Core Ultra 7 258V (Lunar Lake)\nIntel Arc 140V GPU\n32 GB RAM — shared with GPU\nWindows 11',
            font_size=15, color=BODY, margin=0.0)
add_textbox(slide, Inches(0.85), Inches(3.45), Inches(3.3), Inches(0.3), 'Why CPU-only for the model:',
            font_size=16, color=ACCENT, bold=True, margin=0.0)
add_bullets(slide, Inches(0.9), Inches(3.8), Inches(5.25), Inches(1.7), [
    'phi4-mini + 4096 KV cache ≈ 2.6 GB → exceeds Vulkan shared GPU limit → num_gpu=0',
    'Stable at ~5–8 tokens/second — fine for personal use',
], font_size=14.5)
add_code_block(slide, Inches(6.55), Inches(1.8), Inches(5.35), Inches(3.55),
               '# Start all services\n.\\services.ps1 start\n\n# Check health\n.\\services.ps1 health\n\n# Ports:\n# 11434 — Ollama (phi4-mini)\n# 8000  — FastAPI agent\n# 3000  — React UI\n# 8080  — Open WebUI', font_size=12)
add_slide_number(slide, 14)

# Slide 15
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'What I Learned Building This')
lessons_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.55), Inches(11.7), Inches(4.9))
tf = lessons_box.text_frame
tf.clear()
tf.word_wrap = True
lessons = [
    ('Small LLMs: great at NLU, unreliable at multi-step execution', ' — bypass the model for deterministic operations'),
    ('MCP is genuinely useful — and easy to build', ' — 2 days to build, works with Claude Desktop too'),
    ('State management is the hardest part', ' — LangGraph solved persistence; React async state caused subtle bugs'),
    ("Route around LLM limitations, don't fight them", ' — deterministic code beats probabilistic text generation'),
    ('Local AI is viable for personal automation', ' — offline, private, zero cost'),
]
for idx, (lead, desc) in enumerate(lessons, start=1):
    p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(10)
    num_run = p.add_run()
    num_run.text = f'{idx}. '
    style_run(num_run, size=16, color=ACCENT, bold=True)
    lead_run = p.add_run()
    lead_run.text = lead
    style_run(lead_run, size=16, color=TITLE, bold=True)
    desc_run = p.add_run()
    desc_run.text = desc
    style_run(desc_run, size=16, color=BODY)
add_slide_number(slide, 15)

# Slide 16
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'Where This Goes Next')
add_textbox(slide, Inches(1.0), Inches(1.7), Inches(2.0), Inches(0.3), 'Near-term',
            font_size=16, color=ACCENT, bold=True, margin=0.0)
add_bullets(slide, Inches(1.0), Inches(2.05), Inches(5.0), Inches(3.0), [
    'Replace phi4-mini with phi4 (14B) for better tool-calling',
    'refresh_token.txt workflow for 90-day seamless auth',
    'Applied roles tracker to prevent duplicates',
    'Role recommendation scoring based on skills match',
], font_size=15)
add_textbox(slide, Inches(7.0), Inches(1.7), Inches(2.0), Inches(0.3), 'Future',
            font_size=16, color=ACCENT, bold=True, margin=0.0)
add_bullets(slide, Inches(7.0), Inches(2.05), Inches(5.0), Inches(3.0), [
    'Multi-agent: separate search agent and apply agent',
    'Voice interface (Whisper STT + piper TTS)',
    'Calendar integration for start dates / follow-ups',
    'Expose as its own MCP server for Claude Desktop / Copilot',
], font_size=15)
add_slide_number(slide, 16)

# Slide 17
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, "Let's See It Work")
demo = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(2.0), Inches(9.3), Inches(2.8))
apply_fill(demo, CODE_BG)
apply_line(demo, ACCENT, 2)
add_textbox(slide, Inches(2.35), Inches(2.35), Inches(8.6), Inches(2.1),
            '1. "Find data engineer roles in USA"\n   → agent calls search_roles → markdown table\n\n2. "Apply"\n   → bypass: view_role → apply_role (no LLM)\n   → SSE status events stream in real time\n   → ✅ Application submitted\n\n3. (optional) Paste JWT token → ✅ Token updated',
            font_name='Courier New', font_size=13, color=CODE_TEXT)
add_slide_number(slide, 17)

# Slide 18
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, "It's All Open Source")
add_bullets(slide, Inches(0.8), Inches(1.7), Inches(5.25), Inches(3.2), [
    'GitHub: deepakungumaraj/local-llm-phi4',
    'phi4-agent/server.py — FastAPI + bypass layer + SSE',
    'phi4-agent/agent.py — LangGraph agent',
    'phi4-agent/instructions.md — System prompt',
    'services.ps1 — Start/stop all services',
    'myschedule-mcp — MCP server (TypeScript)',
], font_size=15)
add_code_block(slide, Inches(6.5), Inches(1.8), Inches(5.2), Inches(3.4),
               '# server.py — confirmation bypass\nif _is_confirmation(message):\n    apply_params = _get_pending_apply(\n        app.state.agent, thread_id\n    )\n    if apply_params:\n        return EventSourceResponse(\n            _apply_directly_generator(\n                apply_params, thread_id,\n                app.state.mcp_tool_map\n            )\n        )\n    # Path 2: extract role ID\n    role_id = _extract_role_id_from_messages(\n        app.state.agent, thread_id\n    )', font_size=12)
add_slide_number(slide, 18)

# Slide 19
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'Summary')
add_table(slide, Inches(0.7), Inches(1.55), Inches(11.95), Inches(3.1),
          ['Component', 'What It Does', 'Tech'], [
              ['phi4-mini + Ollama', 'Local LLM — understands intent', 'Python, Ollama'],
              ['LangGraph', 'Think→act→observe loop + memory', 'Python'],
              ['MCP + myschedule-mcp', 'Standard protocol + MySchedule tools', 'TypeScript'],
              ['FastAPI + bypass layer', 'Streaming + guaranteed apply execution', 'Python'],
              ['React UI', 'Real-time chat with SSE', 'JavaScript'],
          ], [Inches(3.0), Inches(6.0), Inches(2.95)], font_size=12)
add_callout(slide, Inches(1.5), Inches(5.15), Inches(10.2), Inches(1.05),
            '"AI isn\'t magic. It\'s an LLM, a graph, some tools,\nand a few hundred lines of Python.\nYou can build this too."',
            font_size=15, italic=False)
add_slide_number(slide, 19)

# Slide 20
slide = prs.slides.add_slide(BLANK)
set_bg(slide)
add_title(slide, 'Questions?')
add_textbox(slide, Inches(3.8), Inches(2.0), Inches(5.8), Inches(0.35), 'Deepa Chandramohan',
            font_size=20, color=TITLE, bold=True, align=PP_ALIGN.CENTER, margin=0.0)
add_textbox(slide, Inches(3.2), Inches(2.45), Inches(7.0), Inches(0.3), 'deepa.chandramohan@accenture.com',
            font_size=16, color=ACCENT, align=PP_ALIGN.CENTER, margin=0.0)
add_textbox(slide, Inches(3.0), Inches(2.85), Inches(7.4), Inches(0.3), 'GitHub: deepakungumaraj/local-llm-phi4',
            font_size=14, color=GREY, align=PP_ALIGN.CENTER, margin=0.0)
add_code_block(slide, Inches(2.0), Inches(4.2), Inches(9.3), Inches(1.55), chat_text, font_size=12, color=ACCENT)
add_slide_number(slide, 20)

output_path = r'c:\dev\local-llm-phi4\slide-deck.pptx'
prs.save(output_path)
print(f'Saved presentation to {output_path}')
