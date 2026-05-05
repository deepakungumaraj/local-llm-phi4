from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme
DARK_BG = RGBColor(30, 32, 48)      # #1e2030
ACCENT = RGBColor(0, 212, 255)       # #00d4ff
WHITE = RGBColor(255, 255, 255)
LIGHT_GREY = RGBColor(200, 200, 200)
CODE_BG = RGBColor(20, 22, 35)
HIGHLIGHT_ACCENT = RGBColor(255, 100, 100)

def set_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, subtitle=None):
    """Add title and optional subtitle to slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = text
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.8))
        sub_frame = sub_box.text_frame
        sub_frame.text = subtitle
        sub_frame.word_wrap = True
        p = sub_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE

def add_content_text(slide, text, left=0.5, top=2.2, width=9, height=4.5, size=18):
    """Add body text content"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = text
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = WHITE
        paragraph.space_before = Pt(8)
        paragraph.space_after = Pt(8)
    return text_box

def add_bullet_points(slide, bullets, left=0.5, top=2.2, width=9, height=4.5):
    """Add bullet points to slide"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return text_box

def add_code_block(slide, code, left=0.5, top=3, width=9, height=2.5):
    """Add code block to slide"""
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = CODE_BG
    code_box.line.color.rgb = ACCENT
    code_box.line.width = Pt(1)
    
    text_frame = code_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_bottom = Inches(0.1)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    
    text_frame.text = code
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Courier New"
        paragraph.font.size = Pt(11)
        paragraph.font.color.rgb = ACCENT
    
    return code_box

def add_insight_box(slide, text, left=0.5, top=5.5, width=9, height=1.2):
    """Add key insight callout box"""
    insight_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RGBColor(60, 65, 90)
    insight_box.line.color.rgb = HIGHLIGHT_ACCENT
    insight_box.line.width = Pt(2)
    
    text_frame = insight_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    text_frame.text = text
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = ACCENT
        paragraph.font.italic = True
    
    return insight_box

def add_slide_number(slide, num):
    """Add slide number to bottom right"""
    number_box = slide.shapes.add_textbox(Inches(9.2), Inches(6.8), Inches(0.5), Inches(0.3))
    text_frame = number_box.text_frame
    text_frame.text = str(num)
    p = text_frame.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GREY
    p.alignment = PP_ALIGN.RIGHT

def add_table(slide, rows, cols, left=0.5, top=2.2):
    """Add a table to the slide"""
    width = Inches(9)
    height = Inches(3.5)
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), width, height)
    table = table_shape.table
    
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            
            if row_idx == 0:
                cell.fill.fore_color.rgb = ACCENT
                text_color = DARK_BG
            else:
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(40, 45, 65)
                else:
                    cell.fill.fore_color.rgb = RGBColor(50, 55, 75)
                text_color = WHITE
            
            text_frame = cell.text_frame
            text_frame.word_wrap = True
            if text_frame.paragraphs:
                p = text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.color.rgb = text_color
    
    return table

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Helper to create blank slide
def blank_slide_layout():
    blank_slide_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_slide_layout)

# SLIDE 1: Title Slide
slide1 = blank_slide_layout()
set_background(slide1, DARK_BG)
add_title(slide1, "From Prompt to Production")
add_content_text(slide1, "Building a Personal AI Agent with LangGraph, MCP & Local LLMs", 
                 top=1.8, size=28)
add_content_text(slide1, "Deepa Chandramohan", top=4.5, size=22)
add_content_text(slide1, "AI & Data Practitioner, Accenture", top=5.1, size=18)
add_slide_number(slide1, 1)

# SLIDE 2: The Problem
slide2 = blank_slide_layout()
set_background(slide2, DARK_BG)
add_title(slide2, "The Problem with Manual Staffing")
bullets = [
    "• Internal staffing portal lists hundreds of project roles",
    "• Finding the right role requires repeated manual searching",
    "• Applying requires: search → view → confirm → email → audit → record",
    "• ~10 minutes per application, repeated daily",
    "• Need: automated search + apply flow"
]
add_bullet_points(slide2, bullets)
add_slide_number(slide2, 2)

# SLIDE 3: The Idea
slide3 = blank_slide_layout()
set_background(slide3, DARK_BG)
add_title(slide3, "What If AI Could Do It?")
bullets = [
    "• Goal: type 'find data engineer roles in USA' → get a table",
    "• Type 'apply' → AI searches, views details, submits application",
    "• No cloud APIs. Everything runs locally on my laptop.",
    "• Privacy: resume, profile, auth tokens never leave your machine",
    "• Learn: can small local LLMs handle real multi-step tool use?"
]
add_bullet_points(slide3, bullets)
add_slide_number(slide3, 3)

# SLIDE 4: The Stack
slide4 = blank_slide_layout()
set_background(slide4, DARK_BG)
add_title(slide4, "The Technology Stack")
bullets = [
    "• Language Model: Microsoft phi4-mini (3.8B params) — understands intent",
    "• LLM Runtime: Ollama — serves locally on port 11434",
    "• Agent Orchestration: LangGraph — stateful think→act→observe loop",
    "• AI Framework: LangChain — connects LLM, tools, memory",
    "• Tool Protocol: Model Context Protocol (MCP) — standard tool interface",
    "• Tool Server: Custom Node.js/TypeScript MCP server",
    "• Backend: FastAPI (Python) — HTTP API + SSE streaming",
    "• Frontend: React — chat interface with real-time updates"
]
add_bullet_points(slide4, bullets, height=5)
add_slide_number(slide4, 4)

# SLIDE 5: What is MCP?
slide5 = blank_slide_layout()
set_background(slide5, DARK_BG)
add_title(slide5, "Model Context Protocol (MCP)")
add_title(slide5, "USB-C for AI Tools")
bullets = [
    "• Open standard (Anthropic, November 2024)",
    "• Defines how AI models connect to external tools and data sources",
    "• Before MCP: every AI app needed custom tool integration",
    "• With MCP: build tool server once, any MCP client can use it",
    "• Used by: Claude Desktop, GitHub Copilot, Cursor, and now custom agents"
]
add_bullet_points(slide5, bullets, top=1.8)
add_insight_box(slide5, 
    "USB-C let any device connect to any charger. MCP lets any AI model connect to any tool.",
    top=5.2, height=1)
add_slide_number(slide5, 5)

# SLIDE 6: MCP Server
slide6 = blank_slide_layout()
set_background(slide6, DARK_BG)
add_title(slide6, "MCP Server: Wrapping MySchedule as AI Tools")
add_content_text(slide6, "What it is:", top=2, size=16, height=0.3)
bullets_left = [
    "• Node.js / TypeScript process",
    "• Communicates via stdio pipes (standard input/output)",
    "• Uses @modelcontextprotocol/sdk",
    "• Authenticates with Azure AD OAuth2 tokens"
]
add_bullet_points(slide6, bullets_left, width=4.2, height=2, top=2.3)

add_content_text(slide6, "Tools it exposes:", left=5.2, top=2, size=16, height=0.3)
bullets_right = [
    "• search_roles — keyword/location search",
    "• view_role — full details + projectKey",
    "• apply_role — 5-step apply sequence",
    "• seed_token / seed_refresh_token — auth"
]
add_bullet_points(slide6, bullets_right, left=5.2, width=4.3, height=2, top=2.3)
add_slide_number(slide6, 6)

# SLIDE 7: What is LangGraph?
slide7 = blank_slide_layout()
set_background(slide7, DARK_BG)
add_title(slide7, "LangGraph — Stateful Agent Workflows")
bullets = [
    "• Library from LangChain for building multi-step AI agents",
    "• Define logic as a directed graph: nodes=functions, edges=routing",
    "• Built-in MemorySaver: conversations are stateful across messages",
    "• The loop: think → act → observe → think again (until done)"
]
add_bullet_points(slide7, bullets, height=2)

add_content_text(slide7, "Graph structure:", top=3.8, size=14, height=0.3)
code = """START
  ↓
[agent node] ←─────────────┐
  │ tool call?              │
  ├─ YES → [tools node] ────┘
  │        (MCP call)
  └─ NO → [reporter] → END"""
add_code_block(slide7, code, top=4.1, height=2.2)
add_slide_number(slide7, 7)

# SLIDE 8: Full Architecture
slide8 = blank_slide_layout()
set_background(slide8, DARK_BG)
add_title(slide8, "How Everything Connects")
code = """Browser (React UI)  ──→ FastAPI Server (port 8000)
                           ↓
                    [Bypass Layer]
                    - JWT intercept
                    - Confirmation bypass
                           ↓
                    [LangGraph Agent]
                    - MemorySaver
                    - phi4-mini via Ollama
                           ↓ tool calls
                    [MCP stdio transport]
                           ↓
         [myschedule-mcp Node.js/TypeScript]
         [search_roles / view_role / apply_role]
                           ↓ Azure AD
                    [MySchedule REST API]"""
add_code_block(slide8, code, top=1.8, height=4.8, width=9)
add_slide_number(slide8, 8)

# SLIDE 9: Small LLMs Problem
slide9 = blank_slide_layout()
set_background(slide9, DARK_BG)
add_title(slide9, "The Problem with Small Models")

# Table: What we needed vs What happened
table = add_table(slide9, 3, 2, left=0.5, top=2)
table.cell(0, 0).text_frame.text = "What We Needed"
table.cell(0, 1).text_frame.text = "What phi4-mini Did"
table.cell(1, 0).text_frame.text = "Search → show table"
table.cell(1, 1).text_frame.text = "✅ Worked reliably"
table.cell(2, 0).text_frame.text = "User says 'apply' → call view_role"
table.cell(2, 1).text_frame.text = "❌ Skipped view_role, hallucinated keys"

for row_idx in range(3):
    for col_idx in range(2):
        cell = table.cell(row_idx, col_idx)
        if row_idx == 0:
            cell.fill.fore_color.rgb = ACCENT
            text_color = DARK_BG
        else:
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(40, 45, 65)
            else:
                cell.fill.fore_color.rgb = RGBColor(50, 55, 75)
            text_color = WHITE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = text_color

add_insight_box(slide9,
    "A 3.8B model can't reliably orchestrate 5-step flows. Solution: take LLM out of critical path.",
    top=5.2, height=1)
add_slide_number(slide9, 9)

# SLIDE 10: The Bypass Layer
slide10 = blank_slide_layout()
set_background(slide10, DARK_BG)
add_title(slide10, "The Solution: Server-Side Bypass")
bullets = [
    "• Intercept confirmation messages ('yes', 'confirm') server-side",
    "• Path 1: Check MemorySaver for last view_role → extract projectKey → call apply_role directly",
    "• Path 2: Fallback — regex-scan AI messages for role ID → call view_role → apply_role",
    "• Bonus: JWT token interception — save to disk, model never sees it",
    "• Result: guaranteed correct execution for multi-step flows"
]
add_bullet_points(slide10, bullets)
add_insight_box(slide10,
    "Deterministic code beats probabilistic text generation for multi-step workflows.",
    top=5.2, height=1)
add_slide_number(slide10, 10)

# SLIDE 11: Apply Flow
slide11 = blank_slide_layout()
set_background(slide11, DARK_BG)
add_title(slide11, "Inside apply_role — 5 Steps in One Tool Call")
code = """Step 1: sendEmail ──→ Application email to staffing manager
Step 2: applyToRoleAudit ──→ Record in audit log
Step 3: createCandidate ──→ Create candidate record
Step 4: saveCandidateSelfInput ──→ Save self-input data
Step 5: candidateIndicatorLogic ──→ Log match indicator

Result: ✅ All passed → "Application submitted"
        ⚠️  Partial success → shows step breakdown
        ❌ Failed → validation or auth error"""
add_code_block(slide11, code, top=2, height=4)
add_slide_number(slide11, 11)

# SLIDE 12: Authentication Design
slide12 = blank_slide_layout()
set_background(slide12, DARK_BG)
add_title(slide12, "Authentication Without Storing Secrets")
bullets = [
    "• Access token (token.txt): ~1 hour validity, pasted from browser DevTools",
    "• Refresh token (refresh_token.txt): ~90 days validity, from localStorage",
    "• In-chat token paste: bypass intercepts → re-seeds without model seeing it",
    "• Security: tokens in .gitignore, never committed, never leave machine",
    "• All auth is local — tokens never shared externally"
]
add_bullet_points(slide12, bullets)
add_slide_number(slide12, 12)

# SLIDE 13: SSE Streaming UI
slide13 = blank_slide_layout()
set_background(slide13, DARK_BG)
add_title(slide13, "Real-Time Streaming with Server-Sent Events")
bullets = [
    "• FastAPI uses EventSourceResponse (SSE) to stream agent progress",
    "• Events: status (tool calls), token (text chunks), done (thread_id)",
    "• React UI listens on EventSource, appends tokens in real-time",
    "• Thread persistence: each conversation has thread_id → MemorySaver",
    "• Critical fix: React setState is async — capture thread_id as local var before state update"
]
add_bullet_points(slide13, bullets)
code = """→ [status] "Searching..."
→ [status] "Calling view_role..."
→ [token]  "✅ Application submitted..."
→ [done]   { thread_id: "abc-123" }"""
add_code_block(slide13, code, top=5.2, height=1.8)
add_slide_number(slide13, 13)

# SLIDE 14: Local-First Setup
slide14 = blank_slide_layout()
set_background(slide14, DARK_BG)
add_title(slide14, "Local-First: Everything on Your Laptop")
bullets = [
    "• Hardware: Intel Core Ultra 7 258V, Intel Arc 140V GPU, 32 GB RAM",
    "• Why CPU-only: phi4-mini (2.4GB) + KV cache exceeds Vulkan GPU memory limit",
    "• Performance: ~5–8 tokens/second on CPU (acceptable for personal assistant)",
    "• Service management: services.ps1 start/stop/health commands",
    "• Ports: Ollama 11434, Agent API 8000, React UI 3000, Open WebUI 8080"
]
add_bullet_points(slide14, bullets)
add_slide_number(slide14, 14)

# SLIDE 15: Lessons Learned
slide15 = blank_slide_layout()
set_background(slide15, DARK_BG)
add_title(slide15, "What I Learned Building This")
bullets = [
    "1. Small LLMs are great at NLU, unreliable at multi-step execution",
    "2. MCP is genuinely useful — took 2 days to build, works with any MCP client",
    "3. State management is the hardest part (MemorySaver solved persistence)",
    "4. Never fight LLM limitations — route around them with deterministic code",
    "5. Local AI is viable for personal automation (privacy + cost benefits)"
]
add_bullet_points(slide15, bullets, height=4.5)
add_slide_number(slide15, 15)

# SLIDE 16: What's Next
slide16 = blank_slide_layout()
set_background(slide16, DARK_BG)
add_title(slide16, "Where This Goes Next")
add_content_text(slide16, "Near-term:", top=2, size=16, height=0.3)
bullets_near = [
    "• Replace phi4-mini with phi4 (14B) for better tool-calling",
    "• Add refresh_token.txt workflow for seamless 90-day auth",
    "• Build 'applied roles' tracker to prevent duplicates",
    "• Add role recommendation scoring based on skills match"
]
add_bullet_points(slide16, bullets_near, width=9, height=1.8, top=2.3)

add_content_text(slide16, "Future possibilities:", top=4.3, size=16, height=0.3)
bullets_future = [
    "• Multi-agent: separate search and apply agents",
    "• Voice interface (Whisper STT, piper-tts)",
    "• Calendar integration for start date tracking",
    "• Expose as MCP server — let Claude/Copilot use it"
]
add_bullet_points(slide16, bullets_future, width=9, height=1.6, top=4.6)
add_slide_number(slide16, 16)

# SLIDE 17: Live Demo
slide17 = blank_slide_layout()
set_background(slide17, DARK_BG)
add_title(slide17, "Let's See It Work")
demo_text = """Demo flow:
1. Open React UI at http://localhost:3000
2. Type: "Find data engineer roles in USA"
   Show: agent calls search_roles, returns markdown table
3. Type: "Apply"
   Show: bypass intercepts, calls view_role, then apply_role
   Show: SSE status events streaming in real-time
   Show: ✅ Application submitted with step-by-step results
4. (Optional) Show server logs demonstrating bypass path ran, not LLM

Fallback if auth expired:
   Paste fresh JWT token → bypass intercepts → re-seeds auth → ✅ Token updated"""
add_content_text(slide17, demo_text, top=1.8, size=14, height=5)
add_slide_number(slide17, 17)

# SLIDE 18: The Code
slide18 = blank_slide_layout()
set_background(slide18, DARK_BG)
add_title(slide18, "It's All Open Source")
bullets = [
    "• GitHub repo: deepakungumaraj/local-llm-phi4",
    "• Key files:",
    "  - phi4-agent/server.py — FastAPI + bypass layer + SSE",
    "  - phi4-agent/agent.py — LangGraph agent definition",
    "  - phi4-agent/instructions.md — system prompt",
    "  - services.ps1 — service management",
    "• MCP server: myschedule-mcp (separate TypeScript repo)"
]
add_bullet_points(slide18, bullets, height=4.5)
code = """if _is_confirmation(message):
    apply_params = _get_pending_apply(agent, thread_id)
    if apply_params:
        return EventSourceResponse(
            _apply_directly_generator(apply_params, thread_id, tool_map)
        )"""
add_code_block(slide18, code, top=5.5, height=1.5)
add_slide_number(slide18, 18)

# SLIDE 19: Summary
slide19 = blank_slide_layout()
set_background(slide19, DARK_BG)
add_title(slide19, "Summary")
bullets = [
    "• phi4-mini + Ollama: local LLM, understands intent, decides what to do",
    "• LangGraph: orchestrates think→act→observe loop with memory",
    "• MCP + myschedule-mcp: standard protocol + custom server for tools",
    "• FastAPI + bypass: streaming + guaranteed correct apply execution",
    "• React UI: real-time chat with SSE streaming"
]
add_bullet_points(slide19, bullets, height=3.5)
add_insight_box(slide19,
    "AI isn't magic. It's an LLM, a graph, some tools, and a few hundred lines of Python. You can build this too.",
    top=5.2, height=1.2)
add_slide_number(slide19, 19)

# SLIDE 20: Thank You
slide20 = blank_slide_layout()
set_background(slide20, DARK_BG)
add_title(slide20, "Questions?")
add_content_text(slide20, 
    "Deepa Chandramohan\ndeepa.chandramohan@accenture.com\n\nGitHub: deepakungumaraj/local-llm-phi4",
    top=3, size=20)
add_slide_number(slide20, 20)

# Save presentation
prs.save('slide-deck-opus.pptx')
print("✅ Presentation created: slide-deck-opus.pptx")
print(f"📊 Total slides: {len(prs.slides)}")
